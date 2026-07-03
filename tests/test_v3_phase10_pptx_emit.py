"""Phase 10: editable pptx writer + text merge + OCR annotation contracts."""
from __future__ import annotations

import numpy as np
import pytest
from PIL import Image, ImageDraw

from image_to_editable_ppt.shared.geometry import BBox, ImageSize, Point
from image_to_editable_ppt.v3.compose.merge_text import merge_text_layer_into_slide_ir
from image_to_editable_ppt.v3.core.enums import (
    ConnectorKind,
    ContainerKind,
    NodeKind,
    PortOwnerKind,
    TextRegionRole,
)
from image_to_editable_ppt.v3.emit.models import (
    EmitConnectorPrimitive,
    EmitScene,
    EmitShapePrimitive,
    EmitTextPrimitive,
)
from image_to_editable_ppt.v3.emit.pptx_writer import EMU_PER_PIXEL, write_pptx
from image_to_editable_ppt.v3.emit.style import ShapeVisualStyle, sample_shape_styles
from image_to_editable_ppt.v3.ir.models import (
    PrimitiveNode,
    PrimitiveScene,
    SlideIR,
    TextLayerResult,
    TextRegion,
)


def _scene() -> EmitScene:
    return EmitScene(
        image_size=ImageSize(width=400, height=300),
        shapes=(
            EmitShapePrimitive(
                id="container:1",
                owner_kind=PortOwnerKind.CONTAINER,
                shape_kind=ContainerKind.PANEL,
                bbox=BBox(10, 10, 380, 280),
                confidence=0.9,
            ),
            EmitShapePrimitive(
                id="node:1",
                owner_kind=PortOwnerKind.NODE,
                shape_kind=NodeKind.BOX,
                bbox=BBox(40, 40, 160, 100),
                confidence=0.9,
            ),
            EmitShapePrimitive(
                id="node:2",
                owner_kind=PortOwnerKind.NODE,
                shape_kind=NodeKind.LABEL_ANCHOR,
                bbox=BBox(200, 60, 280, 84),
                confidence=0.9,
            ),
        ),
        texts=(
            EmitTextPrimitive(
                id="text:1",
                role=TextRegionRole.LABEL,
                bbox=BBox(52, 60, 148, 80),
                confidence=0.9,
                text="Encoder",
                owner_ids=("node:1",),
            ),
            EmitTextPrimitive(
                id="text:2",
                role=TextRegionRole.LABEL,
                bbox=BBox(202, 62, 278, 82),
                confidence=0.9,
                text="NP",
                owner_ids=("node:2",),
            ),
            EmitTextPrimitive(
                id="text:3",
                role=TextRegionRole.TITLE,
                bbox=BBox(20, 240, 200, 262),
                confidence=0.9,
                text="Figure title",
            ),
            EmitTextPrimitive(
                id="text:4",
                role=TextRegionRole.LABEL,
                bbox=BBox(300, 240, 360, 260),
                confidence=0.9,
                text=None,  # unrecognized -> must not be emitted
            ),
        ),
        connectors=(
            EmitConnectorPrimitive(
                id="connector:1",
                kind=ConnectorKind.ARROW,
                confidence=0.9,
                source_owner_id="node:1",
                source_owner_kind=PortOwnerKind.NODE,
                target_owner_id="node:2",
                target_owner_kind=PortOwnerKind.NODE,
                source_port_id="port:1",
                target_port_id="port:2",
                path_points=(Point(160, 70), Point(200, 70)),
                arrowhead_end=True,
            ),
            EmitConnectorPrimitive(
                id="connector:2",
                kind=ConnectorKind.LINE,
                confidence=0.9,
                source_owner_id="node:1",
                source_owner_kind=PortOwnerKind.NODE,
                target_owner_id="node:2",
                target_owner_kind=PortOwnerKind.NODE,
                source_port_id="port:3",
                target_port_id="port:4",
                path_points=(),  # no geometry -> must not be drawn
            ),
        ),
    )


def test_write_pptx_writes_native_editable_primitives(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    output = tmp_path / "scene.pptx"
    styles = {"node:1": ShapeVisualStyle(fill=(59, 130, 246), stroke=(30, 64, 175))}
    write_pptx(_scene(), output, styles=styles)

    presentation = Presentation(str(output))
    assert presentation.slide_width == 400 * EMU_PER_PIXEL
    assert presentation.slide_height == 300 * EMU_PER_PIXEL
    shapes = {shape.name: shape for shape in presentation.slides[0].shapes}

    # container + box node are autoshapes; label-anchor node and title are text boxes
    assert shapes["container:1"].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert shapes["node:1"].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert shapes["node:1"].text_frame.text == "Encoder"
    assert str(shapes["node:1"].fill.fore_color.rgb) == "3B82F6"
    assert shapes["node:2"].shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    assert shapes["node:2"].text_frame.text == "NP"
    assert shapes["text:3"].text_frame.text == "Figure title"

    # owned texts are absorbed into shapes, unrecognized standalone text is skipped
    assert "text:1" not in shapes
    assert "text:2" not in shapes
    assert "text:4" not in shapes

    # exactly one connector with an end arrowhead; the geometry-less one is skipped
    assert "connector:2" not in shapes
    connector = shapes["connector:1"]
    assert connector.shape_type == MSO_SHAPE_TYPE.LINE
    assert "tailEnd" in connector._element.xml
    # geometry matches the path in EMU
    assert connector.begin_x == 160 * EMU_PER_PIXEL
    assert connector.end_x == 200 * EMU_PER_PIXEL


def test_write_pptx_positions_match_image_space(tmp_path):
    from pptx import Presentation

    output = tmp_path / "scene.pptx"
    write_pptx(_scene(), output)
    shapes = {shape.name: shape for shape in Presentation(str(output)).slides[0].shapes}
    node = shapes["node:1"]
    assert node.left == 40 * EMU_PER_PIXEL
    assert node.top == 40 * EMU_PER_PIXEL
    assert node.width == 120 * EMU_PER_PIXEL
    assert node.height == 60 * EMU_PER_PIXEL


def test_sample_shape_styles_is_conservative():
    rgb = np.full((300, 400, 3), 255, dtype=np.uint8)
    # filled blue node with dark border
    rgb[40:100, 40:160] = (37, 99, 235)
    scene = _scene()
    styles = sample_shape_styles(rgb, scene)

    node_style = styles["node:1"]
    assert node_style.fill is not None
    assert abs(node_style.fill[2] - 235) <= 3  # blue interior sampled

    # container over plain white background: no fill invented, but stays visible
    container_style = styles["container:1"]
    assert container_style.fill is None
    assert container_style.stroke is not None

    # label anchors are borderless text: never filled/stroked
    label_style = styles["node:2"]
    assert label_style.fill is None and label_style.stroke is None


def test_merge_text_layer_promotes_labels_and_keeps_standalone():
    image_size = ImageSize(width=400, height=300)
    scene = PrimitiveScene(
        image_size=image_size,
        nodes=(
            PrimitiveNode(id="node:1", kind=NodeKind.BOX, bbox=BBox(40, 40, 160, 100), confidence=0.9),
            PrimitiveNode(id="node:2", kind=NodeKind.BOX, bbox=BBox(30, 30, 300, 200), confidence=0.9),
        ),
    )
    slide_ir = SlideIR(image_size=image_size, primitive_scene=scene)
    text_layer = TextLayerResult(
        image_size=image_size,
        regions=(
            TextRegion(id="text:1", bbox=BBox(50, 60, 150, 80), confidence=0.9, text="Encoder"),
            TextRegion(id="text:2", bbox=BBox(200, 250, 280, 270), confidence=0.9, text="caption"),
            TextRegion(id="text:3", bbox=BBox(60, 82, 140, 96), confidence=0.9, text=None),
        ),
    )

    merged = merge_text_layer_into_slide_ir(slide_ir, text_layer)
    merged_scene = merged.primitive_scene
    nodes = {node.id: node for node in merged_scene.nodes}
    # the smallest containing node owns the text and gets the label
    assert nodes["node:1"].label == "Encoder"
    assert nodes["node:2"].label is None
    texts = {text.id: text for text in merged_scene.texts}
    assert texts["text:1"].owner_ids == ("node:1",)
    assert texts["text:2"].owner_ids == ()  # standalone caption
    assert texts["text:3"].owner_ids == ("node:1",)  # owned but unrecognized: no label impact
    assert merged.text_regions == text_layer.regions
    assert merged.text_layer is text_layer


def test_ocr_annotate_fills_confident_text():
    from image_to_editable_ppt.v3.text.ocr import annotate_text_regions, ocr_available

    if not ocr_available():
        pytest.skip("no OCR backend installed")

    image = Image.new("RGB", (240, 60), "white")
    draw = ImageDraw.Draw(image)
    draw.text((12, 20), "Encoder Layer", fill="black")
    regions = (TextRegion(id="text:1", bbox=BBox(8, 14, 120, 36), confidence=0.8),)
    annotated = annotate_text_regions(np.asarray(image, dtype=np.uint8), regions)
    assert annotated[0].text is not None
    assert "encoder" in annotated[0].text.lower()
    assert any(item.startswith("ocr:") for item in annotated[0].provenance)


def test_convert_cli_heuristic_smoke(tmp_path):
    from pptx import Presentation

    from image_to_editable_ppt.ml.convert_to_pptx import main

    image = Image.new("RGB", (320, 240), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((40, 60, 140, 120), outline="black", width=3)
    draw.rectangle((190, 60, 290, 120), outline="black", width=3)
    draw.line((140, 90, 190, 90), fill="black", width=3)
    input_path = tmp_path / "figure.png"
    image.save(input_path)

    output_path = tmp_path / "figure_out.pptx"
    code = main([str(input_path), "-o", str(output_path), "--no-ml", "--no-ocr", "--no-style", "--width", "0"])
    assert code == 0
    assert output_path.exists()
    presentation = Presentation(str(output_path))
    assert presentation.slide_width == 320 * EMU_PER_PIXEL
