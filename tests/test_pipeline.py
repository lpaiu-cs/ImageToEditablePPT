from __future__ import annotations

from pathlib import Path
import zipfile
import xml.etree.ElementTree as ET
import json
from dataclasses import replace

import cv2
import numpy as np
from pptx import Presentation
from PIL import Image, ImageDraw
import pytest

from image_to_editable_ppt.config import PipelineConfig
from image_to_editable_ppt.components import find_connected_components
from image_to_editable_ppt.exporter import export_to_pptx
from image_to_editable_ppt.fitter import (
    fit_component_box_from_outer_contour,
    fit_fill_region_boxes,
    fit_hough_segment_elements,
    hough_axis_strokes,
    hough_connector_elements,
    merge_parallel_strokes,
    Stroke,
)
from image_to_editable_ppt.ir import BBox, BoxGeometry, Element, FillStyle, Point, PolylineGeometry, StrokeStyle
from image_to_editable_ppt.pipeline import build_elements, convert_image
from image_to_editable_ppt.preprocess import ScaleContext, preprocess_image
from image_to_editable_ppt.repair import repair_elements
from image_to_editable_ppt.text import OCRBackend, OCRTextRegion
from image_to_editable_ppt.validation import run_validation_iteration
from image_to_editable_ppt.filtering import filter_residual_components
from image_to_editable_ppt.vlm_parser import DiagramStructure, VLMEdge, VLMNode
from tests.synthetic import (
    boxed_text_cluster_diagram,
    complex_diagram,
    directional_arrow,
    icon_only,
    occluded_box,
    open_contour,
    paper_like_directional_arrow,
    paper_like_dense_text_diagram,
    paper_like_insufficient_widening,
    paper_like_line_with_attached_label_blob,
    paper_like_mixed_arrow_with_connector,
    paper_like_mixed_figure,
    paper_like_multisegment_connector,
    paper_like_noisy_line_ending,
    paper_like_noisy_open_contour,
    paper_like_outer_contour_box_with_label,
    paper_like_occluded_box,
    paper_like_filled_panel_without_border,
    paper_like_symmetric_wedge,
    paper_like_weak_gap_conflict,
    save_image,
    text_box_diagram,
)


class FakeOCRBackend(OCRBackend):
    def __init__(self, regions: list[OCRTextRegion]) -> None:
        self._regions = regions

    def extract(self, image):
        return self._regions


class CropAwareOCRBackend(OCRBackend):
    def __init__(self) -> None:
        self.calls: list[tuple[int, int]] = []

    def extract(self, image):
        self.calls.append(image.size)
        if image.size[0] >= 260:
            return []
        return [
            OCRTextRegion(
                text="Encoder Stage",
                bbox=BBox(8.0, 6.0, min(140.0, image.size[0] - 4.0), min(34.0, image.size[1] - 4.0)),
                confidence=0.97,
            )
        ]


class FakeStructureParser:
    def __init__(self, structure: DiagramStructure) -> None:
        self.structure = structure
        self.calls: list[tuple[int, int]] = []

    def extract_structure(self, image, *, image_path=None):
        self.calls.append(image.size)
        return self.structure


XML_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def assert_arrow_tip_points_to_direction(element: Element, direction: str) -> None:
    start, end = element.geometry.points
    if direction == "right":
        assert end.x > start.x
        return
    if direction == "left":
        assert end.x < start.x
        return
    if direction == "up":
        assert end.y < start.y
        return
    if direction == "down":
        assert end.y > start.y
        return
    raise ValueError(f"unsupported direction: {direction}")


def assert_slide_connector_uses_tail_end_only(output_path: Path, *, expected_tail_count: int = 1) -> None:
    with zipfile.ZipFile(output_path) as archive:
        slide_xml = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
    assert len(slide_xml.findall(".//a:tailEnd", XML_NS)) == expected_tail_count
    assert slide_xml.findall(".//a:headEnd", XML_NS) == []


def test_acceptance_pipeline_detects_core_primitives_and_exports(tmp_path: Path) -> None:
    image_path = save_image(complex_diagram(), tmp_path / "diagram.png")
    output_path = tmp_path / "diagram.pptx"
    result = convert_image(image_path, output_path, config=PipelineConfig())

    kinds = {element.kind for element in result.elements}
    assert "rect" in kinds
    assert "rounded_rect" in kinds
    assert "line" in kinds
    assert "orthogonal_connector" in kinds
    assert "arrow" in kinds
    assert output_path.exists()

    presentation = Presentation(str(output_path))
    slide = presentation.slides[0]
    assert len(slide.shapes) >= len(result.elements)
    for element in result.elements:
        if element.kind in {"rect", "rounded_rect"}:
            if element.fill.enabled:
                assert element.fill.color is not None
        else:
            assert not element.fill.enabled


def test_occlusion_repair_reconstructs_box_with_inferred_flag() -> None:
    result = build_elements(occluded_box(), config=PipelineConfig())
    boxes = [element for element in result.elements if element.kind in {"rect", "rounded_rect"}]
    assert len(boxes) == 1
    assert boxes[0].inferred


def test_open_contour_never_receives_fill() -> None:
    result = build_elements(open_contour(), config=PipelineConfig())
    filled_boxes = [
        element
        for element in result.elements
        if element.kind in {"rect", "rounded_rect"} and element.fill.enabled
    ]
    assert filled_boxes == []


def test_non_diagram_icon_is_omitted() -> None:
    result = build_elements(icon_only(), config=PipelineConfig())
    assert result.elements == []


def test_text_is_only_included_when_high_confidence_and_structural() -> None:
    backend = FakeOCRBackend(
        [
            OCRTextRegion(text="Encoder", bbox=BBox(60.0, 65.0, 150.0, 95.0), confidence=0.97),
            OCRTextRegion(text="noise", bbox=BBox(10.0, 10.0, 30.0, 20.0), confidence=0.98),
            OCRTextRegion(text="low", bbox=BBox(70.0, 70.0, 120.0, 90.0), confidence=0.40),
        ]
    )
    result = build_elements(text_box_diagram(), config=PipelineConfig(), ocr_backend=backend)
    texts = [element for element in result.elements if element.kind == "text"]
    assert len(texts) == 1
    assert texts[0].text is not None
    assert texts[0].text.content == "Encoder"


def test_ocr_routes_text_like_clusters_to_candidate_crops() -> None:
    backend = CropAwareOCRBackend()
    result = build_elements(boxed_text_cluster_diagram(), config=PipelineConfig(), ocr_backend=backend)
    texts = [element for element in result.elements if element.kind == "text"]
    assert texts
    assert texts[0].text is not None
    assert texts[0].text.content == "Encoder Stage"
    assert any(size[0] < 260 for size in backend.calls)


def test_realistic_occluded_box_is_repaired_only_when_evidence_is_strong() -> None:
    result = build_elements(paper_like_occluded_box(), config=PipelineConfig())
    inferred_boxes = [
        element
        for element in result.elements
        if element.kind in {"rect", "rounded_rect"} and element.inferred
    ]
    assert len(inferred_boxes) == 1
    assert inferred_boxes[0].fill.enabled


def test_weak_geometry_with_conflict_does_not_trigger_repair() -> None:
    result = build_elements(paper_like_weak_gap_conflict(), config=PipelineConfig())
    lines = [element for element in result.elements if element.kind == "line"]
    assert lines
    assert all(not element.inferred for element in lines)
    assert max(element.bbox.width for element in lines) < 110


def test_multisegment_orthogonal_connector_is_detected_conservatively() -> None:
    result = build_elements(paper_like_multisegment_connector(), config=PipelineConfig())
    connectors = [element for element in result.elements if element.kind == "orthogonal_connector"]
    assert connectors
    assert any(len(element.geometry.points) >= 4 for element in connectors)


def test_global_segment_graph_recovers_text_occluded_connector_path() -> None:
    mask = np.zeros((220, 340), dtype=bool)
    mask[54:61, 38:138] = True
    mask[54:122, 136:143] = True
    mask[116:123, 136:186] = True
    mask[116:123, 230:288] = True
    mask[116:168, 282:289] = True
    bridge_mask = np.zeros_like(mask)
    bridge_mask[96:146, 182:232] = True
    array = np.full((220, 340, 3), 255, dtype=np.uint8)
    array[mask] = 18
    gray = np.full((220, 340), 255, dtype=np.float32)
    gray[mask] = 18.0
    scale = ScaleContext(
        estimated_stroke_width=3.0,
        min_component_area=18,
        min_stroke_length=16,
        min_linear_length=42,
        min_box_size=24,
    )
    elements = fit_hough_segment_elements(
        mask=mask,
        array=array,
        gray=gray,
        bridge_mask=bridge_mask,
        config=PipelineConfig(),
        scale=scale,
        structural_elements=[],
        existing_elements=[],
        start_index=1,
    )
    connectors = [element for element in elements if element.kind == "orthogonal_connector"]
    assert connectors
    assert any(len(element.geometry.points) >= 4 for element in connectors)


@pytest.mark.parametrize("direction", ["right", "left", "up", "down"])
def test_arrow_detection_orders_points_toward_tip(direction: str) -> None:
    result = build_elements(directional_arrow(direction), config=PipelineConfig())
    arrows = [element for element in result.elements if element.kind == "arrow"]
    assert len(arrows) == 1
    assert_arrow_tip_points_to_direction(arrows[0], direction)


def test_mixed_figure_omits_non_diagram_region() -> None:
    result = build_elements(paper_like_mixed_figure(), config=PipelineConfig())
    assert any(element.kind in {"rect", "line"} for element in result.elements)
    assert all(element.bbox.center.x < 185 for element in result.elements)


def test_no_fill_on_open_contour_under_noise() -> None:
    result = build_elements(paper_like_noisy_open_contour(), config=PipelineConfig())
    assert not any(
        element.kind in {"rect", "rounded_rect"} and element.fill.enabled
        for element in result.elements
    )


def test_dense_text_heavy_diagram_suppresses_text_fragments_with_ocr_off() -> None:
    result = build_elements(paper_like_dense_text_diagram(), config=PipelineConfig())
    large_boxes = [
        element
        for element in result.elements
        if element.kind in {"rect", "rounded_rect"} and element.bbox.width > 260 and element.bbox.height > 150
    ]
    long_connectors = [
        element
        for element in result.elements
        if element.kind in {"line", "orthogonal_connector", "arrow"} and max(element.bbox.width, element.bbox.height) > 150
    ]
    tiny_primitives = [
        element
        for element in result.elements
        if element.kind in {"rect", "rounded_rect", "line", "orthogonal_connector", "arrow"}
        and max(element.bbox.width, element.bbox.height) < 60
    ]
    assert len(large_boxes) >= 2
    assert long_connectors
    assert tiny_primitives == []
    assert not any(element.kind == "text" for element in result.elements)
    assert len(result.elements) <= 8
    assert any(region.reason == "rejected_as_text_like" for region in result.rejected_regions)


def test_unknown_component_is_kept_as_weak_proposal_until_line_fitting() -> None:
    image = paper_like_line_with_attached_label_blob()
    config = PipelineConfig()
    processed = preprocess_image(
        image,
        foreground_threshold=config.foreground_threshold,
        min_component_area=config.min_component_area,
        min_stroke_length=config.min_stroke_length,
        min_box_size=config.min_box_size,
        min_relative_line_length=config.min_relative_line_length,
        min_relative_box_size=config.min_relative_box_size,
        adaptive_background=config.adaptive_background,
        background_blur_divisor=config.background_blur_divisor,
    )
    filtered = filter_residual_components(
        processed.detail_mask,
        processed=processed,
        config=config,
        structural_elements=[],
    )
    assert filtered.diagram_components == []
    assert filtered.weak_components
    assert filtered.rejected_regions == []

    result = build_elements(image, config=config)
    assert any(element.kind == "line" for element in result.elements)


def test_outer_contour_fallback_recovers_box_from_text_merged_component() -> None:
    image = paper_like_outer_contour_box_with_label()
    config = PipelineConfig()
    processed = preprocess_image(
        image,
        foreground_threshold=config.foreground_threshold,
        min_component_area=config.min_component_area,
        min_stroke_length=config.min_stroke_length,
        min_box_size=config.min_box_size,
        min_relative_line_length=config.min_relative_line_length,
        min_relative_box_size=config.min_relative_box_size,
        adaptive_background=config.adaptive_background,
        background_blur_divisor=config.background_blur_divisor,
    )
    components = [
        component
        for component in find_connected_components(processed.detail_mask)
        if component.bbox.width > 180 and component.bbox.height > 120
    ]
    assert components
    component = max(components, key=lambda candidate: candidate.area)
    element = fit_component_box_from_outer_contour(
        component.pixels,
        bbox=component.bbox,
        boundary_mask=processed.boundary_mask_raw,
        array=processed.array,
        detail_mask=processed.detail_mask,
        background_color=processed.background_color,
        config=config,
        scale=processed.scale,
        element_id="box-test",
    )
    assert element is not None
    assert element.kind in {"rect", "rounded_rect"}
    assert element.bbox.width > 220
    assert element.bbox.height > 140


def test_text_cluster_closing_groups_glyphs_into_bridge_blocks() -> None:
    image = boxed_text_cluster_diagram()
    config = PipelineConfig()
    processed = preprocess_image(
        image,
        foreground_threshold=config.foreground_threshold,
        min_component_area=config.min_component_area,
        min_stroke_length=config.min_stroke_length,
        min_box_size=config.min_box_size,
        min_relative_line_length=config.min_relative_line_length,
        min_relative_box_size=config.min_relative_box_size,
        adaptive_background=config.adaptive_background,
        background_blur_divisor=config.background_blur_divisor,
    )
    filtered = filter_residual_components(
        processed.detail_mask_raw,
        processed=processed,
        config=config,
        structural_elements=[],
    )
    assert filtered.text_regions
    assert any(region.width >= 75 and region.height >= 30 for region in filtered.text_regions)


def test_endpoint_to_box_snapping_extends_connector_to_box_edges() -> None:
    image = complex_diagram()
    config = PipelineConfig()
    processed = preprocess_image(
        image,
        foreground_threshold=config.foreground_threshold,
        min_component_area=config.min_component_area,
        min_stroke_length=config.min_stroke_length,
        min_box_size=config.min_box_size,
        min_relative_line_length=config.min_relative_line_length,
        min_relative_box_size=config.min_relative_box_size,
        adaptive_background=config.adaptive_background,
        background_blur_divisor=config.background_blur_divisor,
    )
    boxes = [
        Element(
            id="box-1",
            kind="rect",
            geometry=BoxGeometry(BBox(40.0, 60.0, 120.0, 120.0)),
            stroke=StrokeStyle(color=(0, 0, 0), width=3.0),
            fill=FillStyle(enabled=False, color=None),
            text=None,
            confidence=0.95,
            source_region=BBox(40.0, 60.0, 120.0, 120.0),
        ),
        Element(
            id="box-2",
            kind="rect",
            geometry=BoxGeometry(BBox(210.0, 60.0, 290.0, 120.0)),
            stroke=StrokeStyle(color=(0, 0, 0), width=3.0),
            fill=FillStyle(enabled=False, color=None),
            text=None,
            confidence=0.95,
            source_region=BBox(210.0, 60.0, 290.0, 120.0),
        ),
    ]
    line = Element(
        id="line-1",
        kind="line",
        geometry=PolylineGeometry(points=(Point(128.0, 90.0), Point(202.0, 90.0))),
        stroke=StrokeStyle(color=(0, 0, 0), width=3.0),
        fill=FillStyle(enabled=False, color=None),
        text=None,
        confidence=0.82,
        source_region=BBox(128.0, 88.0, 202.0, 92.0),
    )
    repaired = repair_elements(boxes + [line], processed, config)
    snapped = next(element for element in repaired if element.id == "line-1")
    assert snapped.inferred
    assert snapped.geometry.points[0].x == pytest.approx(120.0)
    assert snapped.geometry.points[-1].x == pytest.approx(210.0)


def test_text_bridge_force_merge_ignores_large_gap_when_bridge_block_dominates() -> None:
    image = complex_diagram()
    config = PipelineConfig()
    processed = preprocess_image(
        image,
        foreground_threshold=config.foreground_threshold,
        min_component_area=config.min_component_area,
        min_stroke_length=config.min_stroke_length,
        min_box_size=config.min_box_size,
        min_relative_line_length=config.min_relative_line_length,
        min_relative_box_size=config.min_relative_box_size,
        adaptive_background=config.adaptive_background,
        background_blur_divisor=config.background_blur_divisor,
    )
    line_a = Element(
        id="line-a",
        kind="line",
        geometry=PolylineGeometry(points=(Point(60.0, 92.0), Point(118.0, 92.0))),
        stroke=StrokeStyle(color=(0, 0, 0), width=3.0),
        fill=FillStyle(enabled=False, color=None),
        text=None,
        confidence=0.84,
        source_region=BBox(60.0, 90.0, 118.0, 94.0),
    )
    line_b = Element(
        id="line-b",
        kind="line",
        geometry=PolylineGeometry(points=(Point(196.0, 92.0), Point(254.0, 92.0))),
        stroke=StrokeStyle(color=(0, 0, 0), width=3.0),
        fill=FillStyle(enabled=False, color=None),
        text=None,
        confidence=0.84,
        source_region=BBox(196.0, 90.0, 254.0, 94.0),
    )
    bridge_mask = np.zeros_like(processed.detail_mask, dtype=bool)
    bridge_mask[74:112, 118:196] = True
    repaired = repair_elements([line_a, line_b], processed, config, bridge_mask=bridge_mask)
    merged_lines = [element for element in repaired if element.kind == "line"]
    assert len(merged_lines) == 1
    assert merged_lines[0].inferred
    assert merged_lines[0].bbox.width >= 190


def test_fill_region_fallback_recovers_borderless_panel() -> None:
    result = build_elements(paper_like_filled_panel_without_border(), config=PipelineConfig())
    boxes = [element for element in result.elements if element.kind in {"rect", "rounded_rect"}]
    assert boxes
    assert any(box.fill.enabled and box.bbox.width > 240 and box.bbox.height > 140 for box in boxes)


def test_preprocess_blacklists_textured_photo_region_for_graph_input() -> None:
    image = Image.new("RGB", (280, 190), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((24, 36, 114, 118), outline="black", width=5)
    rng = np.random.default_rng(91)
    patch = rng.integers(0, 255, size=(110, 104, 3), dtype=np.uint8)
    image_array = np.asarray(image).copy()
    image_array[34:144, 152:256, :] = patch
    image = Image.fromarray(image_array)
    config = PipelineConfig()
    processed = preprocess_image(
        image,
        foreground_threshold=config.foreground_threshold,
        min_component_area=config.min_component_area,
        min_stroke_length=config.min_stroke_length,
        min_box_size=config.min_box_size,
        min_relative_line_length=config.min_relative_line_length,
        min_relative_box_size=config.min_relative_box_size,
        adaptive_background=config.adaptive_background,
        background_blur_divisor=config.background_blur_divisor,
        fill_region_background_ratio=config.fill_region_background_ratio,
        fill_region_uniformity_ratio=config.fill_region_uniformity_ratio,
        fill_region_edge_ratio=config.fill_region_edge_ratio,
        non_diagram_edge_density=config.non_diagram_edge_density,
        non_diagram_color_variance=config.non_diagram_color_variance,
        non_diagram_side_support=config.non_diagram_side_support,
    )

    assert float(processed.non_diagram_mask[34:144, 152:256].mean()) >= 0.72
    assert float(processed.non_diagram_mask[24:124, 24:124].mean()) <= 0.08
    assert float(processed.boundary_mask[34:144, 152:256].mean()) <= float(processed.boundary_mask_raw[34:144, 152:256].mean()) * 0.2


def test_global_segment_graph_rejects_floating_chain_when_box_anchored_path_exists() -> None:
    mask = np.zeros((220, 320), dtype=bool)
    array = np.full((220, 320, 3), 255, dtype=np.uint8)

    def draw_stroke(stroke: Stroke) -> None:
        x0 = max(0, int(np.floor(stroke.x0)))
        y0 = max(0, int(np.floor(stroke.y0)))
        x1 = min(mask.shape[1], int(np.ceil(stroke.x1)))
        y1 = min(mask.shape[0], int(np.ceil(stroke.y1)))
        mask[y0:y1, x0:x1] = True
        array[y0:y1, x0:x1, :] = 22

    strokes = [
        Stroke("horizontal", 70.0, 78.5, 126.0, 81.5, 3.0),
        Stroke("vertical", 124.5, 80.0, 127.5, 160.0, 3.0),
        Stroke("horizontal", 126.0, 158.5, 200.0, 161.5, 3.0),
        Stroke("horizontal", 24.0, 18.5, 154.0, 21.5, 3.0),
        Stroke("vertical", 152.5, 20.0, 155.5, 92.0, 3.0),
        Stroke("horizontal", 154.0, 89.5, 300.0, 92.5, 3.0),
    ]
    for stroke in strokes:
        draw_stroke(stroke)

    scale = ScaleContext(
        estimated_stroke_width=3.0,
        min_component_area=18,
        min_stroke_length=16,
        min_linear_length=42,
        min_box_size=24,
    )
    boxes = [
        Element(
            id="box-left",
            kind="rect",
            geometry=BoxGeometry(BBox(20.0, 50.0, 70.0, 110.0)),
            stroke=StrokeStyle(color=(0, 0, 0), width=3.0),
            fill=FillStyle(enabled=False, color=None),
            text=None,
            confidence=0.95,
            source_region=BBox(20.0, 50.0, 70.0, 110.0),
        ),
        Element(
            id="box-right",
            kind="rect",
            geometry=BoxGeometry(BBox(200.0, 130.0, 252.0, 190.0)),
            stroke=StrokeStyle(color=(0, 0, 0), width=3.0),
            fill=FillStyle(enabled=False, color=None),
            text=None,
            confidence=0.95,
            source_region=BBox(200.0, 130.0, 252.0, 190.0),
        ),
    ]
    elements, _ = hough_connector_elements(
        strokes,
        mask=mask,
        array=array,
        config=PipelineConfig(),
        scale=scale,
        bridge_mask=None,
        structural_elements=boxes,
        existing_elements=[],
        start_index=1,
    )

    assert any(element.kind == "orthogonal_connector" and element.bbox.x0 <= 70.0 and element.bbox.x1 >= 200.0 for element in elements)
    assert all(element.bbox.y0 >= 60.0 for element in elements)


def test_fill_region_hierarchy_splits_nested_panel_from_parent_contour() -> None:
    mask = np.zeros((220, 220), dtype=bool)
    mask[20:190, 20:200] = True
    mask[78:146, 84:152] = False
    boundary_mask = np.zeros_like(mask)
    boundary_mask[76:80, 84:152] = True
    boundary_mask[144:148, 84:152] = True
    boundary_mask[78:146, 82:86] = True
    boundary_mask[78:146, 150:154] = True
    array = np.full((220, 220, 3), 255, dtype=np.uint8)
    array[20:190, 20:200, :] = (224, 233, 244)
    array[78:146, 84:152, :] = (230, 238, 246)
    smoothed = array.copy()
    detail_mask = np.zeros_like(mask)
    scale = ScaleContext(
        estimated_stroke_width=3.0,
        min_component_area=18,
        min_stroke_length=16,
        min_linear_length=42,
        min_box_size=24,
    )

    boxes = fit_fill_region_boxes(
        mask=mask,
        boundary_mask=boundary_mask,
        array=array,
        smoothed_array=smoothed,
        detail_mask=detail_mask,
        background_color=(255, 255, 255),
        config=PipelineConfig(),
        scale=scale,
        existing_elements=[],
        start_index=1,
    )

    assert any(box.bbox.width >= 160 and box.bbox.height >= 140 for box in boxes)
    assert any(60 <= box.bbox.x0 <= 95 and 70 <= box.bbox.y0 <= 90 and box.bbox.width >= 55 and box.bbox.height >= 55 for box in boxes)


def test_hough_bridge_mask_merges_segments_across_text_gap() -> None:
    mask = np.zeros((140, 320), dtype=bool)
    mask[64:73, 26:150] = True
    mask[64:73, 190:294] = True
    bridge_mask = np.zeros_like(mask)
    bridge_mask[46:96, 150:190] = True
    array = np.full((140, 320, 3), 255, dtype=np.uint8)
    array[mask] = 18
    gray = np.full((140, 320), 255, dtype=np.float32)
    gray[mask] = 18.0
    scale = ScaleContext(
        estimated_stroke_width=3.0,
        min_component_area=18,
        min_stroke_length=16,
        min_linear_length=42,
        min_box_size=24,
    )
    detected = cv2.HoughLinesP(
        (mask.astype(np.uint8) * 255),
        1.0,
        np.pi / 180.0,
        threshold=16,
        minLineLength=14,
        maxLineGap=12,
    )
    assert detected is not None
    raw_strokes = hough_axis_strokes(detected, mask=mask, scale=scale, min_length=14)
    config = replace(PipelineConfig(), stroke_merge_gap=72)
    merged = merge_parallel_strokes(
        raw_strokes,
        config,
        mask=mask,
        array=array,
        gray=gray,
        allow_gap_merge=True,
        bridge_mask=bridge_mask,
    )
    assert any(stroke.orientation == "horizontal" and stroke.length >= 240 and stroke.inferred for stroke in merged)


def test_arrow_exporter_unit_maps_tip_to_ooxml_tail_end(tmp_path: Path) -> None:
    output_path = tmp_path / "arrow.pptx"
    export_to_pptx(
        [
            Element(
                id="arrow-1",
                kind="arrow",
                geometry=PolylineGeometry(points=(Point(20.0, 40.0), Point(180.0, 40.0))),
                stroke=StrokeStyle(color=(0, 0, 0), width=4.0),
                fill=FillStyle(enabled=False, color=None),
                text=None,
                confidence=0.95,
                source_region=BBox(20.0, 34.0, 180.0, 46.0),
            )
        ],
        (200, 80),
        output_path,
        PipelineConfig(),
    )
    assert_slide_connector_uses_tail_end_only(output_path)


@pytest.mark.parametrize("direction", ["right", "left", "up", "down"])
def test_arrow_convert_image_exports_tip_to_ooxml_tail_end(
    tmp_path: Path,
    direction: str,
) -> None:
    image_path = save_image(directional_arrow(direction), tmp_path / f"{direction}-arrow.png")
    output_path = tmp_path / f"{direction}-arrow.pptx"
    result = convert_image(image_path, output_path, config=PipelineConfig())
    arrows = [element for element in result.elements if element.kind == "arrow"]
    assert len(arrows) == 1
    assert_arrow_tip_points_to_direction(arrows[0], direction)
    assert_slide_connector_uses_tail_end_only(output_path)


@pytest.mark.parametrize("direction", ["right", "left", "up", "down"])
def test_paper_like_arrow_convert_image_exports_tip_to_ooxml_tail_end(
    tmp_path: Path,
    direction: str,
) -> None:
    image_path = save_image(paper_like_directional_arrow(direction), tmp_path / f"paper-like-{direction}-arrow.png")
    output_path = tmp_path / f"paper-like-{direction}-arrow.pptx"
    result = convert_image(image_path, output_path, config=PipelineConfig())
    arrows = [element for element in result.elements if element.kind == "arrow"]
    assert len(arrows) == 1
    assert_arrow_tip_points_to_direction(arrows[0], direction)
    assert_slide_connector_uses_tail_end_only(output_path)


def test_paper_like_insufficient_widening_degrades_to_line(tmp_path: Path) -> None:
    image_path = save_image(paper_like_insufficient_widening(), tmp_path / "insufficient-widening.png")
    output_path = tmp_path / "insufficient-widening.pptx"
    result = convert_image(image_path, output_path, config=PipelineConfig())
    arrows = [element for element in result.elements if element.kind == "arrow"]
    lines = [element for element in result.elements if element.kind == "line"]
    assert arrows == []
    assert lines
    assert_slide_connector_uses_tail_end_only(output_path, expected_tail_count=0)


def test_paper_like_symmetric_wedge_is_omitted(tmp_path: Path) -> None:
    image_path = save_image(paper_like_symmetric_wedge(), tmp_path / "symmetric-wedge.png")
    output_path = tmp_path / "symmetric-wedge.pptx"
    result = convert_image(image_path, output_path, config=PipelineConfig())
    assert result.elements == []
    assert_slide_connector_uses_tail_end_only(output_path, expected_tail_count=0)


def test_paper_like_noisy_line_ending_is_not_exported_as_arrow(tmp_path: Path) -> None:
    image_path = save_image(paper_like_noisy_line_ending(), tmp_path / "noisy-line-ending.png")
    output_path = tmp_path / "noisy-line-ending.pptx"
    result = convert_image(image_path, output_path, config=PipelineConfig())
    arrows = [element for element in result.elements if element.kind == "arrow"]
    lines = [element for element in result.elements if element.kind == "line"]
    assert arrows == []
    assert lines
    assert_slide_connector_uses_tail_end_only(output_path, expected_tail_count=0)


def test_paper_like_mixed_arrow_exports_single_tail_end_marker(tmp_path: Path) -> None:
    image_path = save_image(paper_like_mixed_arrow_with_connector(), tmp_path / "mixed-arrow.png")
    output_path = tmp_path / "mixed-arrow.pptx"
    result = convert_image(image_path, output_path, config=PipelineConfig())
    arrows = [element for element in result.elements if element.kind == "arrow"]
    other_linear = [
        element
        for element in result.elements
        if element.kind in {"line", "orthogonal_connector"}
    ]
    assert len(arrows) == 1
    assert_arrow_tip_points_to_direction(arrows[0], "right")
    assert other_linear
    assert_slide_connector_uses_tail_end_only(output_path)


def test_convert_image_debug_dump_serializes_slot_dataclasses(tmp_path: Path) -> None:
    image_path = save_image(complex_diagram(), tmp_path / "debug-diagram.png")
    output_path = tmp_path / "debug-diagram.pptx"
    debug_path = tmp_path / "debug-elements.json"
    convert_image(image_path, output_path, config=PipelineConfig(), debug_elements_path=debug_path)
    assert debug_path.exists()
    rejection_path = tmp_path / "debug-elements.rejections.json"
    assert rejection_path.exists()
    with debug_path.open("r", encoding="utf-8") as handle:
        data = json.load(handle)
    assert isinstance(data, list)
    assert data
    assert "geometry" in data[0]


def test_semantic_pipeline_uses_vlm_structure_and_local_snapping() -> None:
    image = Image.new("RGB", (320, 200), "white")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((42, 52, 132, 112), radius=16, outline=(24, 78, 128), fill=(220, 235, 248), width=4)
    draw.rounded_rectangle((206, 84, 296, 144), radius=16, outline=(64, 64, 64), fill=(248, 228, 206), width=4)
    parser = FakeStructureParser(
        DiagramStructure(
            nodes=[
                VLMNode("n1", "box", "Vector Store", BBox(88.0, 200.0, 444.0, 610.0)),
                VLMNode("n2", "box", "Planner", BBox(606.0, 360.0, 950.0, 770.0)),
            ],
            edges=[VLMEdge("n1", "n2", "solid_arrow", "retrieves")],
            coordinate_space="normalized_1000",
        )
    )

    result = build_elements(
        image,
        config=PipelineConfig(semantic_fallback_to_legacy=False),
        structure_parser=parser,
    )

    assert result.pipeline_mode == "semantic"
    boxes = [element for element in result.elements if element.kind in {"rect", "rounded_rect"}]
    texts = [element for element in result.elements if element.kind == "text"]
    arrows = [element for element in result.elements if element.kind == "arrow"]
    assert len(boxes) == 2
    assert len(arrows) == 1
    assert any(text.text is not None and text.text.content == "Vector Store" for text in texts)
    assert any(text.text is not None and text.text.content == "Planner" for text in texts)
    assert any(text.text is not None and text.text.content == "retrieves" for text in texts)
    assert boxes[0].bbox.x0 == pytest.approx(42.0, abs=8.0)
    assert boxes[0].bbox.y0 == pytest.approx(52.0, abs=8.0)
    assert len(arrows[0].geometry.points) >= 3


def test_semantic_routed_arrow_exports_single_tail_marker(tmp_path: Path) -> None:
    image = Image.new("RGB", (320, 220), "white")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((36, 48, 126, 108), radius=16, outline="black", fill=(224, 240, 250), width=4)
    draw.rounded_rectangle((206, 126, 296, 186), radius=16, outline="black", fill=(250, 232, 212), width=4)
    image_path = save_image(image, tmp_path / "semantic.png")
    output_path = tmp_path / "semantic.pptx"
    parser = FakeStructureParser(
        DiagramStructure(
            nodes=[
                VLMNode("n1", "box", "Memory", BBox(75.0, 164.0, 431.0, 545.0)),
                VLMNode("n2", "box", "Planner", BBox(606.0, 527.0, 956.0, 882.0)),
            ],
            edges=[VLMEdge("n1", "n2", "dashed_arrow")],
            coordinate_space="normalized_1000",
        )
    )

    result = convert_image(
        image_path,
        output_path,
        config=PipelineConfig(semantic_fallback_to_legacy=False),
        structure_parser=parser,
    )

    assert result.pipeline_mode == "semantic"
    assert any(element.kind == "arrow" and len(element.geometry.points) >= 3 for element in result.elements)
    assert_slide_connector_uses_tail_end_only(output_path)


def test_validation_iteration_exports_pptx_svg_and_comparison_artifacts(tmp_path: Path) -> None:
    image_path = save_image(complex_diagram(), tmp_path / "validate-diagram.png")
    run = run_validation_iteration(image_path, tmp_path / "iter", config=PipelineConfig())
    assert run.artifacts.output_pptx.exists()
    assert run.artifacts.output_svg.exists()
    assert run.artifacts.rendered_png.exists()
    assert run.artifacts.overlay_png.exists()
    assert run.artifacts.edge_diff_png.exists()
    assert run.artifacts.metrics_json.exists()
    assert run.metrics.rendered_shape_count >= 4
    svg_text = run.artifacts.output_svg.read_text(encoding="utf-8")
    assert "<svg" in svg_text
