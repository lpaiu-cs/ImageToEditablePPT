from __future__ import annotations

import json
import random
from pathlib import Path

import pytest
from PIL import Image

import image_to_editable_ppt.ml.generate_dataset as generate_dataset_cli
from image_to_editable_ppt.ml.adapter import AnnotationMLAdapter
from image_to_editable_ppt.ml.annotation_schema import (
    SCHEMA_VERSION,
    AnnotationImageSize,
    DetectorAnnotationDocument,
)
from image_to_editable_ppt.ml.generate_dataset import assign_splits
from image_to_editable_ppt.ml.synthesize import (
    find_soffice,
    generate_slide_spec,
    render_spec_image,
    validate_spec_contract,
    write_spec_pptx,
)
from image_to_editable_ppt.v3.core.enums import DiagramFamily
from image_to_editable_ppt.v3.ir.validate import validate_slide_ir


def test_generated_specs_satisfy_slide_ir_contract() -> None:
    rng = random.Random(42)
    for index in range(25):
        spec = generate_slide_spec(rng, sample_id=f"contract_{index:03d}")
        validate_spec_contract(spec)
        # orthogonal_flow now spans a single-row chain (3-6 nodes) and a realistic
        # 2D-grid block diagram (more nodes, variable connectors), so the per-node
        # text-region pairing is the layout-agnostic invariant to assert here.
        assert 3 <= len(spec.nodes) <= 16
        assert len(spec.connectors) >= 1
        assert len(spec.text_regions) == len(spec.nodes)


def test_decorations_are_rendered_but_not_ground_truth() -> None:
    # Decorations (braces etc.) are opt-in hard negatives: drawn into the image but
    # never GT, and only when explicitly requested (off by default).
    assert generate_slide_spec(random.Random(0), sample_id="d").decorations == ()  # default off
    spec = None
    for seed in range(60):
        candidate = generate_slide_spec(random.Random(seed), sample_id=f"deco_{seed}", with_decorations=True)
        if candidate.decorations:
            spec = candidate
            break
    assert spec is not None, "no decoration produced in 60 seeds"
    document = spec.to_annotation_document()
    assert len(document.primitive_scene.connector_candidates) == len(spec.connectors)
    # rendering with a decoration must not raise (PIL path)
    render_spec_image(spec)


def test_generation_is_deterministic_for_equal_seeds() -> None:
    spec_a = generate_slide_spec(random.Random(123), sample_id="det")
    spec_b = generate_slide_spec(random.Random(123), sample_id="det")

    assert spec_a.to_annotation_document() == spec_b.to_annotation_document()
    assert render_spec_image(spec_a).tobytes() == render_spec_image(spec_b).tobytes()

    spec_c = generate_slide_spec(random.Random(124), sample_id="det")
    assert spec_a.to_annotation_document() != spec_c.to_annotation_document()


def test_generator_rejects_unsupported_family() -> None:
    with pytest.raises(ValueError, match="unsupported synthetic family"):
        generate_slide_spec(random.Random(1), sample_id="bad", family=DiagramFamily.SWIMLANE)


def test_generated_cycle_specs_satisfy_slide_ir_contract() -> None:
    rng = random.Random(42)
    for index in range(25):
        spec = generate_slide_spec(rng, sample_id=f"cycle_{index:03d}", family=DiagramFamily.CYCLE)
        validate_spec_contract(spec)
        assert spec.family is DiagramFamily.CYCLE
        assert 3 <= len(spec.nodes) <= 6
        # a cycle is a closed loop: one connector per node, forming a ring
        assert len(spec.connectors) == len(spec.nodes)
        assert len(spec.text_regions) == len(spec.nodes)
        assert spec.family_proposal.family is DiagramFamily.CYCLE


def test_generated_table_matrix_specs_satisfy_slide_ir_contract() -> None:
    rng = random.Random(44)
    for index in range(25):
        spec = generate_slide_spec(rng, sample_id=f"tm_{index:03d}", family=DiagramFamily.TABLE_MATRIX)
        validate_spec_contract(spec)
        assert spec.family is DiagramFamily.TABLE_MATRIX
        assert len(spec.nodes) >= 4  # at least a 2x2 grid
        assert spec.connectors == ()  # a matrix has no connectors
        assert spec.containers == ()
        assert len(spec.text_regions) == len(spec.nodes)


def test_generated_block_flow_specs_are_trees_with_branching_connectors() -> None:
    rng = random.Random(45)
    saw_branch = False
    for index in range(25):
        spec = generate_slide_spec(rng, sample_id=f"bf_{index:03d}", family=DiagramFamily.BLOCK_FLOW)
        validate_spec_contract(spec)
        assert spec.family is DiagramFamily.BLOCK_FLOW
        assert len(spec.nodes) >= 3
        assert len(spec.connectors) == len(spec.nodes) - 1  # a tree: edges = nodes - 1
        roots = {c.candidate.start_endpoint.owner_id for c in spec.connectors}
        if any(
            sum(1 for c in spec.connectors if c.candidate.start_endpoint.owner_id == r) >= 2 for r in roots
        ):
            saw_branch = True
    assert saw_branch  # at least one slide has a node with multiple out-edges


def test_cycle_connectors_form_a_closed_ring() -> None:
    spec = generate_slide_spec(random.Random(7), sample_id="ring", family=DiagramFamily.CYCLE)
    node_ids = [node.id for node in spec.nodes]
    edges = {
        (connector.candidate.start_endpoint.owner_id, connector.candidate.end_endpoint.owner_id)
        for connector in spec.connectors
    }
    expected = {(node_ids[i], node_ids[(i + 1) % len(node_ids)]) for i in range(len(node_ids))}
    assert edges == expected


def test_render_draws_structures_at_annotated_positions() -> None:
    spec = generate_slide_spec(random.Random(5), sample_id="render")
    image = render_spec_image(spec)

    assert image.size == (spec.image_size.width, spec.image_size.height)
    background = image.getpixel((2, 2))
    assert background == (255, 255, 255)
    # Each node must render *something* within its annotated bbox (fill, outline,
    # or — for text-only/mono styles — the label glyphs). Probing exact centre/edge
    # pixels is too brittle now that nodes can be unfilled, text-only, or elliptical.
    for node in spec.nodes:
        region = image.crop(
            (int(node.bbox.x0), int(node.bbox.y0), int(node.bbox.x1) + 1, int(node.bbox.y1) + 1)
        )
        extrema = region.getextrema()
        assert any(channel_min < 255 for channel_min, _ in extrema), f"node {node.id} region is blank"


def test_pptx_sidecar_preserves_shape_mapping(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")

    spec = generate_slide_spec(random.Random(9), sample_id="pptx")
    pptx_path = tmp_path / "sample.pptx"
    write_spec_pptx(spec, pptx_path)

    presentation = pptx.Presentation(pptx_path)
    assert presentation.slide_width == spec.image_size.width * 9525
    assert presentation.slide_height == spec.image_size.height * 9525
    shape_names = [shape.name for shape in presentation.slides[0].shapes]
    for node in spec.nodes:
        assert node.id in shape_names
    for container in spec.containers:
        assert container.id in shape_names
    for connector in spec.connectors:
        assert connector.candidate.id in shape_names


def test_pptx_container_uses_varied_style_color(tmp_path: Path) -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.dml.color import RGBColor

    spec = next(
        (s for i in range(60) if (s := generate_slide_spec(random.Random(i), sample_id=f"c{i}")).containers),
        None,
    )
    assert spec is not None and spec.container_styles
    path = tmp_path / "container.pptx"
    write_spec_pptx(spec, path)

    shapes = {shape.name: shape for shape in pptx.Presentation(path).slides[0].shapes}
    container = shapes[spec.containers[0].id]
    # The pptx container must honour the per-sample style, not the old faint hardcoded color.
    assert container.fill.fore_color.rgb == RGBColor(*spec.container_styles[0].fill)
    assert container.line.color.rgb == RGBColor(*spec.container_styles[0].outline)


def test_assign_splits_partitions_every_sample() -> None:
    splits = assign_splits(100, rng=random.Random(3), train_ratio=0.8, val_ratio=0.1)

    assert len(splits) == 100
    assert splits.count("train") == 80
    assert splits.count("val") == 10
    assert splits.count("test") == 10
    assert assign_splits(100, rng=random.Random(3), train_ratio=0.8, val_ratio=0.1) == splits


def test_generate_dataset_cli_writes_triplets_and_manifest(tmp_path: Path) -> None:
    output_dir = tmp_path / "dataset"
    exit_code = generate_dataset_cli.main(
        [
            "--output-dir",
            str(output_dir),
            "--count",
            "6",
            "--seed",
            "21",
            "--image-width",
            "640",
            "--image-height",
            "360",
            "--train-ratio",
            "0.5",
            "--val-ratio",
            "0.25",
        ]
    )
    assert exit_code == 0

    manifest = json.loads((output_dir / "dataset_manifest.json").read_text(encoding="utf-8"))
    assert manifest["schema_version"] == SCHEMA_VERSION
    assert manifest["generator"]["seed"] == 21
    assert manifest["split_counts"] == {"train": 3, "val": 2, "test": 1}
    assert sum(manifest["family_counts"].values()) == 6
    assert len(manifest["samples"]) == 6

    adapter = AnnotationMLAdapter()
    for sample in manifest["samples"]:
        image_path = output_dir / sample["image"]
        annotation_path = output_dir / sample["annotation"]
        pptx_path = output_dir / sample["pptx"]
        assert image_path.exists() and annotation_path.exists() and pptx_path.exists()
        assert Path(sample["image"]).parts[0] == sample["split"]

        with Image.open(image_path) as image:
            assert image.size == (640, 360)
        document = DetectorAnnotationDocument.from_dict(json.loads(annotation_path.read_text(encoding="utf-8")))
        assert document.image_id == sample["id"]
        assert document.split == sample["split"]
        assert document.image_size == AnnotationImageSize(width=640, height=360)
        validate_slide_ir(adapter.to_slide_ir(document))


@pytest.mark.skipif(find_soffice() is None, reason="LibreOffice (soffice) not installed")
def test_generate_dataset_cli_supports_soffice_renderer(tmp_path: Path) -> None:
    output_dir = tmp_path / "dataset"
    exit_code = generate_dataset_cli.main(
        ["--output-dir", str(output_dir), "--count", "1", "--seed", "2", "--renderer", "soffice"]
    )
    assert exit_code == 0

    manifest = json.loads((output_dir / "dataset_manifest.json").read_text(encoding="utf-8"))
    assert manifest["generator"]["renderer"] == "soffice_png_v1"
    sample = manifest["samples"][0]
    with Image.open(output_dir / sample["image"]) as image:
        assert image.size == (1280, 720)
    assert (output_dir / sample["pptx"]).exists()


def test_generate_dataset_cli_can_skip_pptx(tmp_path: Path) -> None:
    output_dir = tmp_path / "dataset"
    exit_code = generate_dataset_cli.main(
        ["--output-dir", str(output_dir), "--count", "2", "--seed", "1", "--no-pptx"]
    )
    assert exit_code == 0

    manifest = json.loads((output_dir / "dataset_manifest.json").read_text(encoding="utf-8"))
    assert manifest["generator"]["pptx_written"] is False
    for sample in manifest["samples"]:
        assert sample["pptx"] is None
    assert not list(output_dir.rglob("*.pptx"))
