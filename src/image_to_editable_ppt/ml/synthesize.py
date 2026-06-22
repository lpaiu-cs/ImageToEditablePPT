"""Synthetic dataset generation for the Phase 7 ML track.

Strategy (see plan.md, Phase 7 / Task 2): instead of hand-labelling real
slides, programmatically compose random diagram structures, write them both
as an editable .pptx and as a rendered .png, and emit the ground-truth
``DetectorAnnotationDocument`` from the same spec. The spec is the single
source of truth, so image, pptx, and annotation geometry always agree.

Rendering uses a deterministic PIL rasterizer rather than a PowerPoint
renderer so dataset generation does not depend on LibreOffice/Office being
installed; the .pptx sidecar preserves the img - ppt mapping pair.
"""

from __future__ import annotations

import math
import random
from dataclasses import dataclass, replace
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

from image_to_editable_ppt.ml.adapter import AnnotationMLAdapter, DetectorModelOutput
from image_to_editable_ppt.ml.annotation_schema import (
    AnnotationBBox,
    AnnotationConnectorCandidate,
    AnnotationConnectorEndpoint,
    AnnotationContainer,
    AnnotationFamilyProposal,
    AnnotationImageSize,
    AnnotationNode,
    AnnotationPoint,
    AnnotationPort,
    AnnotationTextRegion,
    DetectorAnnotationDocument,
)
from image_to_editable_ppt.v3.core.enums import (
    ConnectorKind,
    ContainerKind,
    DiagramFamily,
    NodeKind,
    PortOwnerKind,
    PortSide,
    TextRegionRole,
)
from image_to_editable_ppt.v3.ir.validate import validate_slide_ir

GENERATOR_NAME = "synthetic_ppt_render_v1"
RENDERER_NAME = "pil_raster_v1"
EMU_PER_PIXEL = 9525  # 1 px at 96 dpi

# Container visual style for the PIL rasterizer. Containers are drawn with a
# visible, *varied* style picked per-sample (see _pick_container_style) so the
# detector learns container presence from a realistic range of boundaries rather
# than one fixed look. These module constants are the fallback used when a spec
# carries no explicit style; they are deliberately visible (a saturated outline),
# unlike the original near-invisible faint-gray rendering which made the
# container signal unlearnable.
CONTAINER_FILL = (248, 250, 252)
CONTAINER_OUTLINE = (71, 85, 105)
CONTAINER_OUTLINE_WIDTH = 3

# Palettes for per-container styling. Every outline is saturated/dark enough to
# be clearly visible against the white background and the light fills. Fills lean
# toward grey panels because real paper group boxes are most often grey/white
# rounded or dashed panels rather than saturated colours.
_CONTAINER_OUTLINE_PALETTE = (
    (71, 85, 105),    # slate
    (55, 65, 81),     # gray-700
    (31, 41, 55),     # gray-800
    (37, 99, 235),    # blue
    (190, 18, 60),    # rose
    (15, 118, 110),   # teal
    (180, 83, 9),     # amber
    (109, 40, 217),   # violet
)
_CONTAINER_FILL_PALETTE = (
    (255, 255, 255),  # white (border-only container)
    (255, 255, 255),  # white (weighted: many real panels are border-only)
    (248, 250, 252),  # slate-50
    (241, 245, 249),  # slate-100
    (236, 238, 241),  # grey panel
    (228, 230, 235),  # darker grey panel
    (239, 246, 255),  # blue-50
    (240, 253, 244),  # green-50
    (254, 249, 235),  # amber-50
)
_CONTAINER_OUTLINE_WIDTHS = (2, 3, 4)

SUPPORTED_FAMILIES = (
    DiagramFamily.ORTHOGONAL_FLOW,
    DiagramFamily.CYCLE,
    DiagramFamily.TABLE_MATRIX,
    DiagramFamily.BLOCK_FLOW,
    DiagramFamily.GRAPH,
    DiagramFamily.TREE,
    DiagramFamily.LAYERED_STACK,
)

_LABEL_VOCAB = (
    "Plan",
    "Input",
    "Parse",
    "Review",
    "Build",
    "Test",
    "Deploy",
    "Verify",
    "Ship",
    "Audit",
    "Merge",
    "Release",
)

# Short labels for the node-link / stack families (graph states, parse-tree
# constituents, NN layers) — distinct vocab so the families look the part.
_GRAPH_LABEL_VOCAB = (
    "q0", "q1", "q2", "q3", "S", "A", "B", "C", "D", "E", "v1", "v2", "v3", "n1", "n2",
)
_TREE_LABEL_VOCAB = (
    "S", "NP", "VP", "PP", "Det", "N", "V", "Adj", "Adv", "the", "dog", "runs",
    "saw", "cat", "ROOT", "IP", "CP", "DP", "AdvP",
)
_LAYER_LABEL_VOCAB = (
    "Input", "Embed", "Conv", "Pool", "Linear", "ReLU", "Softmax", "LSTM",
    "Attention", "Dropout", "Output", "Encoder", "Decoder", "LayerNorm",
)

_FILL_PALETTE = (
    (219, 234, 254),
    (220, 252, 231),
    (254, 243, 199),
    (252, 231, 243),
    (237, 233, 254),
    (226, 232, 240),
)

_STROKE_PALETTE = (
    (30, 64, 175),
    (22, 101, 52),
    (146, 64, 14),
    (157, 23, 77),
    (91, 33, 182),
    (51, 65, 85),
)

# Grayscale palettes for domain randomization toward the real-paper look. Many
# arXiv/ACL figures are greyscale or black-and-white line art rather than the
# bright filled boxes the original generator produced; training on a mix of
# colour / grey / mono themes closes that synthetic->real appearance gap.
_GRAY_FILL_PALETTE = (
    (255, 255, 255),
    (245, 245, 245),
    (238, 238, 240),
    (229, 231, 235),
    (224, 224, 228),
)
_GRAY_STROKE_PALETTE = (
    (0, 0, 0),
    (31, 41, 55),
    (55, 65, 81),
    (75, 85, 99),
)

# Connector strokes for the thin-black ("paper") theme — real arrows are usually
# a 1-2px black/dark line, unlike the synthetic 3px coloured connectors.
_THIN_CONNECTOR_STROKES = (
    (0, 0, 0),
    (31, 41, 55),
    (51, 65, 85),
)


@dataclass(slots=True, frozen=True)
class RenderTheme:
    """Per-sample visual randomization knobs (domain randomization).

    Chosen once per generated sample so the dataset spans the colour / greyscale
    / black-and-white spectrum of real paper figures and a mix of thick-coloured
    vs thin-black connectors. Only *appearance* is randomized — geometry and the
    emitted ground truth are unaffected.
    """

    mode: str  # "color" | "gray" | "mono"
    thin_connectors: bool
    node_stroke_width: int
    connector_width: int
    connector_stroke: tuple[int, int, int]


def _pick_theme(rng: random.Random) -> RenderTheme:
    mode = rng.choices(("color", "gray", "mono"), weights=(0.4, 0.4, 0.2))[0]
    thin = rng.random() < 0.6  # real arrows skew thin-black
    return RenderTheme(
        mode=mode,
        thin_connectors=thin,
        node_stroke_width=rng.choice((2, 2, 3)),
        connector_width=rng.choice((1, 2)) if thin else 3,
        connector_stroke=_THIN_CONNECTOR_STROKES[rng.randrange(len(_THIN_CONNECTOR_STROKES))]
        if thin
        else (0, 0, 0),  # overridden per-connector to the node stroke when not thin
    )


def _theme_fill_stroke(rng: random.Random, theme: RenderTheme) -> tuple[tuple[int, int, int], tuple[int, int, int]]:
    if theme.mode == "color":
        index = rng.randrange(len(_FILL_PALETTE))
        return _FILL_PALETTE[index], _STROKE_PALETTE[index]
    if theme.mode == "gray":
        return (
            _GRAY_FILL_PALETTE[rng.randrange(len(_GRAY_FILL_PALETTE))],
            _GRAY_STROKE_PALETTE[rng.randrange(len(_GRAY_STROKE_PALETTE))],
        )
    return (255, 255, 255), (0, 0, 0)  # mono


def _connector_stroke(theme: RenderTheme, node_stroke: tuple[int, int, int]) -> tuple[int, int, int]:
    return theme.connector_stroke if theme.thin_connectors else node_stroke


@dataclass(slots=True, frozen=True)
class SyntheticNodeStyle:
    fill: tuple[int, int, int] | None
    stroke: tuple[int, int, int] | None
    shape: str = "auto"  # auto (from NodeKind) | rect | round | ellipse | text
    stroke_width: int = 3


@dataclass(slots=True, frozen=True)
class SyntheticContainerStyle:
    fill: tuple[int, int, int]
    outline: tuple[int, int, int]
    outline_width: int
    dashed: bool = False
    shape: str = "round"  # round | rect


def _pick_container_style(rng: random.Random) -> SyntheticContainerStyle:
    """Pick a visible, varied container style (grey/white panel, dashed or solid).

    Dashed rectangular group boxes and plain grey panels are the dominant real
    paper container looks, so the sampler weights toward them; styling only
    affects appearance, never the emitted ground-truth geometry.
    """
    dashed = rng.random() < 0.4
    return SyntheticContainerStyle(
        fill=_CONTAINER_FILL_PALETTE[rng.randrange(len(_CONTAINER_FILL_PALETTE))],
        outline=_CONTAINER_OUTLINE_PALETTE[rng.randrange(len(_CONTAINER_OUTLINE_PALETTE))],
        outline_width=_CONTAINER_OUTLINE_WIDTHS[rng.randrange(len(_CONTAINER_OUTLINE_WIDTHS))],
        dashed=dashed,
        # dashed group boxes are almost always rectangular; solid ones often rounded.
        shape="rect" if dashed or rng.random() < 0.4 else "round",
    )


@dataclass(slots=True, frozen=True)
class SyntheticConnector:
    candidate: AnnotationConnectorCandidate
    start_port: AnnotationPort
    end_port: AnnotationPort
    stroke: tuple[int, int, int]
    width: int = 3
    draw_arrowhead: bool = True


@dataclass(slots=True, frozen=True)
class SyntheticDecoration:
    """A non-connector mark (brace, bracket, annotation line) drawn as a polyline.

    Rendered into the image but deliberately NOT ground truth — these are the hard
    negatives (the curly brace grouping two output nodes is the canonical real-figure
    false positive). Classical extraction proposes them as connectors; seeing them
    labelled as non-edges is what teaches the relation judge to reject them."""

    points: tuple[tuple[float, float], ...]
    stroke: tuple[int, int, int] = (0, 0, 0)
    width: int = 2


@dataclass(slots=True, frozen=True)
class SyntheticSlideSpec:
    sample_id: str
    family: DiagramFamily
    image_size: AnnotationImageSize
    nodes: tuple[AnnotationNode, ...]
    node_styles: dict[str, SyntheticNodeStyle]
    text_regions: tuple[AnnotationTextRegion, ...]
    containers: tuple[AnnotationContainer, ...]
    connectors: tuple[SyntheticConnector, ...]
    family_proposal: AnnotationFamilyProposal
    label_font_size: int
    container_styles: tuple[SyntheticContainerStyle, ...] = ()
    decorations: tuple[SyntheticDecoration, ...] = ()  # rendered, non-GT hard negatives

    def to_annotation_document(self, *, split: str | None = None, metadata: dict[str, object] | None = None) -> DetectorAnnotationDocument:
        output = DetectorModelOutput(
            image_id=self.sample_id,
            image_size=self.image_size,
            family_predictions=(self.family_proposal,),
            node_predictions=self.nodes,
            container_predictions=self.containers,
            text_predictions=self.text_regions,
            port_predictions=tuple(
                port for connector in self.connectors for port in (connector.start_port, connector.end_port)
            ),
            # Persist the rendered stroke width so the connector segmenter's GT mask
            # matches the rendered line thickness (thin-arrow randomization).
            connector_predictions=tuple(
                replace(connector.candidate, stroke_width=connector.width) for connector in self.connectors
            ),
            metadata={
                "generator": GENERATOR_NAME,
                "renderer": RENDERER_NAME,
                "synthetic_ground_truth": True,
                **(metadata or {}),
            },
        )
        document = _GT_ADAPTER.from_model_output(output)
        if split is None:
            return document
        return DetectorAnnotationDocument(
            image_id=document.image_id,
            image_size=document.image_size,
            schema_version=document.schema_version,
            image_path=document.image_path,
            split=split,
            family_proposals=document.family_proposals,
            text_regions=document.text_regions,
            diagram_instances=document.diagram_instances,
            primitive_scene=document.primitive_scene,
            metadata=document.metadata,
        )


_GT_ADAPTER = AnnotationMLAdapter(default_source="synthetic_gt", default_provenance_prefix="synthetic_gt")


def generate_slide_spec(
    rng: random.Random,
    *,
    sample_id: str,
    family: DiagramFamily = DiagramFamily.ORTHOGONAL_FLOW,
    image_size: AnnotationImageSize | None = None,
) -> SyntheticSlideSpec:
    if family not in SUPPORTED_FAMILIES:
        raise ValueError(f"unsupported synthetic family: {family.value} (supported: {[f.value for f in SUPPORTED_FAMILIES]})")
    size = image_size or AnnotationImageSize(width=1280, height=720)
    if family is DiagramFamily.CYCLE:
        spec = _generate_cycle_spec(rng, sample_id=sample_id, image_size=size)
    elif family is DiagramFamily.TABLE_MATRIX:
        spec = _generate_table_matrix_spec(rng, sample_id=sample_id, image_size=size)
    elif family is DiagramFamily.BLOCK_FLOW:
        spec = _generate_block_flow_spec(rng, sample_id=sample_id, image_size=size)
    elif family is DiagramFamily.GRAPH:
        spec = _generate_graph_spec(rng, sample_id=sample_id, image_size=size)
    elif family is DiagramFamily.TREE:
        spec = _generate_tree_spec(rng, sample_id=sample_id, image_size=size)
    elif family is DiagramFamily.LAYERED_STACK:
        spec = _generate_layered_stack_spec(rng, sample_id=sample_id, image_size=size)
    elif family is DiagramFamily.ORTHOGONAL_FLOW and rng.random() < 0.55:
        # Half of orthogonal flows use the realistic 2D-grid block-diagram layout
        # with solid-black elbow connectors (the dominant real-paper look).
        spec = _generate_grid_flow_spec(rng, sample_id=sample_id, image_size=size)
    else:
        spec = _generate_orthogonal_flow_spec(rng, sample_id=sample_id, image_size=size)
    decorations = _sample_decorations(rng, spec)
    return replace(spec, decorations=decorations) if decorations else spec


def _brace_points(
    x_tip: float, y0: float, y1: float, depth: float, side: int
) -> tuple[tuple[float, float], ...]:
    """A curly-brace ``}``/``{`` polyline: two tips at (x_tip, y0/y1) with the spine
    bulging ``side`` (±1) by ``depth`` and a point poking out twice as far at the mid."""
    ymid = (y0 + y1) / 2.0
    return (
        (x_tip, y0),
        (x_tip + side * depth, y0 + (ymid - y0) * 0.3),
        (x_tip + side * depth, ymid - (ymid - y0) * 0.3),
        (x_tip + side * 2 * depth, ymid),
        (x_tip + side * depth, ymid + (y1 - ymid) * 0.3),
        (x_tip + side * depth, y1 - (y1 - ymid) * 0.3),
        (x_tip, y1),
    )


def _sample_decorations(rng: random.Random, spec: SyntheticSlideSpec) -> tuple[SyntheticDecoration, ...]:
    """Occasionally add a non-GT brace grouping the top- and bottom-most nodes in the
    right margin — the canonical real-figure false connector (cf. the "Résumé" brace
    grouping the two output boxes). It must span two nodes that are NOT truly connected,
    so the classical candidate it induces is a genuine negative, not a mislabelled
    positive on top of a real edge."""
    nodes = spec.nodes
    if len(nodes) < 3 or rng.random() >= 0.45:
        return ()
    idx = {n.id: i for i, n in enumerate(nodes)}
    gt_pairs = set()
    for con in spec.connectors:
        s = con.candidate.start_endpoint.owner_id if con.candidate.start_endpoint else None
        t = con.candidate.end_endpoint.owner_id if con.candidate.end_endpoint else None
        if s in idx and t in idx:
            gt_pairs.add(frozenset((idx[s], idx[t])))
    order = sorted(range(len(nodes)), key=lambda i: (nodes[i].bbox.y0 + nodes[i].bbox.y1) / 2.0)
    top, bot = nodes[order[0]], nodes[order[-1]]
    if frozenset((order[0], order[-1])) in gt_pairs:
        return ()  # would sit on a real edge; skip rather than mislabel
    # Place the brace in the clear margin right of EVERY node and container, so it is a
    # single isolated component (not cut into per-row fragments by node erasure, which
    # would otherwise attach to adjacent connected pairs and mislabel them positive).
    width = float(spec.image_size.width)
    right_limit = max([n.bbox.x1 for n in nodes] + [c.bbox.x1 for c in spec.containers])
    depth = rng.randint(10, 20)
    x_tip = right_limit + rng.randint(24, 50)
    if x_tip + 2 * depth > width - 4 or bot.bbox.y1 - top.bbox.y0 <= 60:
        return ()  # no room in the right margin, or the span is too short to be brace-like
    points = _brace_points(x_tip, top.bbox.y0, bot.bbox.y1, depth, side=1)
    return (SyntheticDecoration(points=points, width=rng.randint(2, 3)),)


def _generate_grid_flow_spec(
    rng: random.Random,
    *,
    sample_id: str,
    image_size: AnnotationImageSize,
) -> SyntheticSlideSpec:
    """A realistic block diagram: a 2D grid of boxes wired by solid-black elbow
    (right-angle) connectors, like the PowerPoint-built architecture figures that
    fill real papers. The connector segmenter trained only on thin straight
    single-row connectors barely fired on these, so this is the key real-transfer
    lever for connectors/topology."""
    theme = _pick_theme(rng)
    fill, stroke = _theme_fill_stroke(rng, theme)
    node_kind = NodeKind.ROUNDED_BOX if rng.random() < 0.5 else NodeKind.BOX
    node_shape = "round" if node_kind is NodeKind.ROUNDED_BOX else "rect"
    style = SyntheticNodeStyle(fill=fill, stroke=stroke, shape=node_shape, stroke_width=theme.node_stroke_width)
    # Real connectors are solid black/dark, medium weight, with triangle heads.
    connector_stroke = (0, 0, 0) if rng.random() < 0.8 else _GRAY_STROKE_PALETTE[rng.randrange(len(_GRAY_STROKE_PALETTE))]
    connector_width = rng.choice((2, 3))

    width, height = float(image_size.width), float(image_size.height)
    rows, cols = rng.randint(2, 4), rng.randint(2, 4)
    margin_x, margin_y = width * 0.07, height * 0.1
    cell_w = (width - 2 * margin_x) / cols
    cell_h = (height - 2 * margin_y) / rows
    box_w = cell_w * rng.uniform(0.55, 0.74)
    box_h = cell_h * rng.uniform(0.4, 0.62)
    label_font_size = max(13, int(box_h * 0.3))
    font = ImageFont.load_default(size=label_font_size)

    nodes: list[AnnotationNode] = []
    text_regions: list[AnnotationTextRegion] = []
    node_styles: dict[str, SyntheticNodeStyle] = {}
    grid: dict[tuple[int, int], int] = {}
    index = 0
    for r in range(rows):
        for c in range(cols):
            if rng.random() < 0.15 and not (r == 0 and c == 0):
                continue  # occasional empty cell, like real diagrams
            cx = margin_x + c * cell_w + cell_w / 2.0
            cy = margin_y + r * cell_h + cell_h / 2.0
            bbox = AnnotationBBox(x0=cx - box_w / 2.0, y0=cy - box_h / 2.0, x1=cx + box_w / 2.0, y1=cy + box_h / 2.0)
            node_id = f"node:{sample_id}:{index}"
            label = rng.choice(_LABEL_VOCAB)
            nodes.append(
                AnnotationNode(
                    id=node_id, kind=node_kind, bbox=bbox, confidence=1.0, label=label,
                    text_region_ids=(f"text:{sample_id}:{index}",), source="synthetic_gt",
                    provenance=(f"{GENERATOR_NAME}:node",),
                )
            )
            node_styles[node_id] = style
            text_regions.append(_make_text_region(label, sample_id=sample_id, index=index, font=font, node_bbox=bbox))
            grid[(r, c)] = index
            index += 1

    connectors: list[SyntheticConnector] = []
    edge_index = 0

    def connect(a: int, b: int, start_side: PortSide, end_side: PortSide) -> None:
        nonlocal edge_index
        start_point = _side_point(nodes[a].bbox, start_side)
        end_point = _side_point(nodes[b].bbox, end_side)
        connectors.append(
            _make_connector(
                sample_id=sample_id, index=edge_index, start_node=nodes[a], end_node=nodes[b],
                start_point=start_point, start_side=start_side, end_point=end_point, end_side=end_side,
                theme=theme, node_stroke=stroke, path=_elbow_path(start_point, start_side, end_point, end_side),
                stroke_override=connector_stroke, width_override=connector_width,
            )
        )
        edge_index += 1

    for (r, c), ni in grid.items():
        if (r, c + 1) in grid and rng.random() < 0.8:
            connect(ni, grid[(r, c + 1)], PortSide.RIGHT, PortSide.LEFT)
        if (r + 1, c) in grid and rng.random() < 0.7:
            connect(ni, grid[(r + 1, c)], PortSide.BOTTOM, PortSide.TOP)

    if not connectors and len(nodes) >= 2:
        # A block diagram is always wired: if the probabilistic pass produced nothing,
        # force the first adjacent grid pair (or any two nodes) so there is >=1 edge.
        for (r, c), ni in grid.items():
            if (r, c + 1) in grid:
                connect(ni, grid[(r, c + 1)], PortSide.RIGHT, PortSide.LEFT)
                break
            if (r + 1, c) in grid:
                connect(ni, grid[(r + 1, c)], PortSide.BOTTOM, PortSide.TOP)
                break
        else:
            connect(0, 1, PortSide.RIGHT, PortSide.LEFT)

    containers, container_styles = _build_containers(
        rng, sample_id=sample_id, nodes=nodes, image_size=image_size, probability=0.5
    )
    focus = containers[0].bbox if containers else _union_bbox(
        [node.bbox for node in nodes] + [connector.candidate.bbox for connector in connectors]
    )
    proposal = AnnotationFamilyProposal(
        id=f"family:{sample_id}:0", family=DiagramFamily.ORTHOGONAL_FLOW, confidence=1.0, focus_bbox=focus,
        evidence=(f"{GENERATOR_NAME}:layout",), provenance=(f"{GENERATOR_NAME}:family_proposal",),
    )
    return SyntheticSlideSpec(
        sample_id=sample_id, family=DiagramFamily.ORTHOGONAL_FLOW, image_size=image_size, nodes=tuple(nodes),
        node_styles=node_styles, text_regions=tuple(text_regions), containers=containers, connectors=tuple(connectors),
        family_proposal=proposal, label_font_size=label_font_size, container_styles=container_styles,
    )


def _generate_orthogonal_flow_spec(
    rng: random.Random,
    *,
    sample_id: str,
    image_size: AnnotationImageSize,
) -> SyntheticSlideSpec:
    horizontal = rng.random() < 0.7
    node_count = rng.randint(3, 6)
    node_kind = NodeKind.ROUNDED_BOX if rng.random() < 0.5 else NodeKind.BOX
    theme = _pick_theme(rng)
    _fill, _stroke = _theme_fill_stroke(rng, theme)
    style = SyntheticNodeStyle(fill=_fill, stroke=_stroke, stroke_width=theme.node_stroke_width)

    width, height = float(image_size.width), float(image_size.height)
    margin_x, margin_y = width * 0.08, height * 0.12
    lane_extent = (width if horizontal else height) - 2 * (margin_x if horizontal else margin_y)
    slot = lane_extent / node_count
    gap = slot * rng.uniform(0.25, 0.4)
    node_extent = slot - gap
    cross_extent = (height if horizontal else width) * rng.uniform(0.14, 0.22)
    cross_center = (height if horizontal else width) * rng.uniform(0.35, 0.65)

    label_font_size = max(14, int(cross_extent * 0.3))
    font = ImageFont.load_default(size=label_font_size)
    labels = rng.sample(_LABEL_VOCAB, node_count)

    nodes: list[AnnotationNode] = []
    text_regions: list[AnnotationTextRegion] = []
    node_styles: dict[str, SyntheticNodeStyle] = {}
    for index in range(node_count):
        lane_start = (margin_x if horizontal else margin_y) + slot * index + gap / 2.0
        jitter = rng.uniform(-gap * 0.15, gap * 0.15)
        if horizontal:
            bbox = AnnotationBBox(
                x0=lane_start,
                y0=cross_center - cross_extent / 2.0 + jitter,
                x1=lane_start + node_extent,
                y1=cross_center + cross_extent / 2.0 + jitter,
            )
        else:
            bbox = AnnotationBBox(
                x0=cross_center - cross_extent / 2.0 + jitter,
                y0=lane_start,
                x1=cross_center + cross_extent / 2.0 + jitter,
                y1=lane_start + node_extent,
            )
        node_id = f"node:{sample_id}:{index}"
        label = labels[index]
        text_id = f"text:{sample_id}:{index}"
        nodes.append(
            AnnotationNode(
                id=node_id,
                kind=node_kind,
                bbox=bbox,
                confidence=1.0,
                label=label,
                text_region_ids=(text_id,),
                source="synthetic_gt",
                provenance=(f"{GENERATOR_NAME}:node",),
            )
        )
        node_styles[node_id] = style
        text_regions.append(
            AnnotationTextRegion(
                id=text_id,
                bbox=_label_bbox(label, font=font, node_bbox=bbox),
                confidence=1.0,
                role=TextRegionRole.LABEL,
                text=label,
                source="synthetic_gt",
                provenance=(f"{GENERATOR_NAME}:text",),
            )
        )

    connectors: list[SyntheticConnector] = []
    for index in range(node_count - 1):
        start_node, end_node = nodes[index], nodes[index + 1]
        if horizontal:
            start_side, end_side = PortSide.RIGHT, PortSide.LEFT
            start_point = AnnotationPoint(start_node.bbox.x1, (start_node.bbox.y0 + start_node.bbox.y1) / 2.0)
            end_point = AnnotationPoint(end_node.bbox.x0, (end_node.bbox.y0 + end_node.bbox.y1) / 2.0)
        else:
            start_side, end_side = PortSide.BOTTOM, PortSide.TOP
            start_point = AnnotationPoint((start_node.bbox.x0 + start_node.bbox.x1) / 2.0, start_node.bbox.y1)
            end_point = AnnotationPoint((end_node.bbox.x0 + end_node.bbox.x1) / 2.0, end_node.bbox.y0)
        connector_id = f"connector:{sample_id}:{index}"
        path = _orthogonal_path(start_point, end_point, horizontal=horizontal)
        connectors.append(
            SyntheticConnector(
                candidate=AnnotationConnectorCandidate(
                    id=connector_id,
                    kind=ConnectorKind.ARROW,
                    bbox=_path_bbox(path),
                    confidence=1.0,
                    source_evidence_id=f"evidence:{connector_id}",
                    path_points=path,
                    start_endpoint=AnnotationConnectorEndpoint(
                        point=start_point,
                        owner_id=start_node.id,
                        owner_kind=PortOwnerKind.NODE,
                        side=start_side,
                    ),
                    end_endpoint=AnnotationConnectorEndpoint(
                        point=end_point,
                        owner_id=end_node.id,
                        owner_kind=PortOwnerKind.NODE,
                        side=end_side,
                    ),
                    arrowhead_end=True,
                    source="synthetic_gt",
                    provenance=(f"{GENERATOR_NAME}:connector",),
                ),
                start_port=_port_for(start_node.id, sample_id, index, "start", side=start_side, point=start_point),
                end_port=_port_for(end_node.id, sample_id, index, "end", side=end_side, point=end_point),
                stroke=_connector_stroke(theme, style.stroke),
                width=theme.connector_width,
            )
        )

    containers, container_styles = _build_containers(
        rng, sample_id=sample_id, nodes=nodes, image_size=image_size, probability=0.65
    )
    focus = containers[0].bbox if containers else _union_bbox(
        [node.bbox for node in nodes] + [connector.candidate.bbox for connector in connectors]
    )
    proposal = AnnotationFamilyProposal(
        id=f"family:{sample_id}:0",
        family=DiagramFamily.ORTHOGONAL_FLOW,
        confidence=1.0,
        focus_bbox=focus,
        evidence=(f"{GENERATOR_NAME}:layout",),
        provenance=(f"{GENERATOR_NAME}:family_proposal",),
    )

    return SyntheticSlideSpec(
        sample_id=sample_id,
        family=DiagramFamily.ORTHOGONAL_FLOW,
        image_size=image_size,
        nodes=tuple(nodes),
        node_styles=node_styles,
        text_regions=tuple(text_regions),
        containers=containers,
        connectors=tuple(connectors),
        family_proposal=proposal,
        label_font_size=label_font_size,
        container_styles=container_styles,
    )


def _generate_cycle_spec(
    rng: random.Random,
    *,
    sample_id: str,
    image_size: AnnotationImageSize,
) -> SyntheticSlideSpec:
    """Nodes evenly placed on a ring, linked into a closed directed loop."""
    node_count = rng.randint(3, 6)
    node_kind = NodeKind.ROUNDED_BOX if rng.random() < 0.5 else NodeKind.BOX
    theme = _pick_theme(rng)
    _fill, _stroke = _theme_fill_stroke(rng, theme)
    style = SyntheticNodeStyle(fill=_fill, stroke=_stroke, stroke_width=theme.node_stroke_width)

    width, height = float(image_size.width), float(image_size.height)
    center_x, center_y = width / 2.0, height / 2.0
    radius_x = width * rng.uniform(0.27, 0.33)
    radius_y = height * rng.uniform(0.27, 0.33)
    node_half_w = width * rng.uniform(0.07, 0.09)
    node_half_h = height * rng.uniform(0.06, 0.08)
    start_angle = -math.pi / 2.0 + rng.uniform(-0.15, 0.15)

    label_font_size = max(14, int(node_half_h * 2.0 * 0.3))
    font = ImageFont.load_default(size=label_font_size)
    labels = rng.sample(_LABEL_VOCAB, node_count)

    nodes: list[AnnotationNode] = []
    text_regions: list[AnnotationTextRegion] = []
    node_styles: dict[str, SyntheticNodeStyle] = {}
    centers: list[AnnotationPoint] = []
    for index in range(node_count):
        angle = start_angle + 2.0 * math.pi * index / node_count
        cx = center_x + radius_x * math.cos(angle)
        cy = center_y + radius_y * math.sin(angle)
        centers.append(AnnotationPoint(cx, cy))
        bbox = AnnotationBBox(
            x0=cx - node_half_w,
            y0=cy - node_half_h,
            x1=cx + node_half_w,
            y1=cy + node_half_h,
        )
        node_id = f"node:{sample_id}:{index}"
        label = labels[index]
        text_id = f"text:{sample_id}:{index}"
        nodes.append(
            AnnotationNode(
                id=node_id,
                kind=node_kind,
                bbox=bbox,
                confidence=1.0,
                label=label,
                text_region_ids=(text_id,),
                source="synthetic_gt",
                provenance=(f"{GENERATOR_NAME}:node",),
            )
        )
        node_styles[node_id] = style
        text_regions.append(
            AnnotationTextRegion(
                id=text_id,
                bbox=_label_bbox(label, font=font, node_bbox=bbox),
                confidence=1.0,
                role=TextRegionRole.LABEL,
                text=label,
                source="synthetic_gt",
                provenance=(f"{GENERATOR_NAME}:text",),
            )
        )

    connectors: list[SyntheticConnector] = []
    for index in range(node_count):
        start_node = nodes[index]
        end_node = nodes[(index + 1) % node_count]
        start_point, start_side = _edge_point_toward(start_node.bbox, centers[(index + 1) % node_count])
        end_point, end_side = _edge_point_toward(end_node.bbox, centers[index])
        connector_id = f"connector:{sample_id}:{index}"
        path = (start_point, end_point)
        connectors.append(
            SyntheticConnector(
                candidate=AnnotationConnectorCandidate(
                    id=connector_id,
                    kind=ConnectorKind.ARROW,
                    bbox=_path_bbox(path),
                    confidence=1.0,
                    source_evidence_id=f"evidence:{connector_id}",
                    path_points=path,
                    start_endpoint=AnnotationConnectorEndpoint(
                        point=start_point,
                        owner_id=start_node.id,
                        owner_kind=PortOwnerKind.NODE,
                        side=start_side,
                    ),
                    end_endpoint=AnnotationConnectorEndpoint(
                        point=end_point,
                        owner_id=end_node.id,
                        owner_kind=PortOwnerKind.NODE,
                        side=end_side,
                    ),
                    arrowhead_end=True,
                    source="synthetic_gt",
                    provenance=(f"{GENERATOR_NAME}:connector",),
                ),
                start_port=_port_for(start_node.id, sample_id, index, "start", side=start_side, point=start_point),
                end_port=_port_for(end_node.id, sample_id, index, "end", side=end_side, point=end_point),
                stroke=_connector_stroke(theme, style.stroke),
                width=theme.connector_width,
            )
        )

    containers, container_styles = _build_containers(
        rng, sample_id=sample_id, nodes=nodes, image_size=image_size, probability=0.5
    )
    focus = containers[0].bbox if containers else _union_bbox(
        [node.bbox for node in nodes] + [connector.candidate.bbox for connector in connectors]
    )
    proposal = AnnotationFamilyProposal(
        id=f"family:{sample_id}:0",
        family=DiagramFamily.CYCLE,
        confidence=1.0,
        focus_bbox=focus,
        evidence=(f"{GENERATOR_NAME}:layout",),
        provenance=(f"{GENERATOR_NAME}:family_proposal",),
    )

    return SyntheticSlideSpec(
        sample_id=sample_id,
        family=DiagramFamily.CYCLE,
        image_size=image_size,
        nodes=tuple(nodes),
        node_styles=node_styles,
        text_regions=tuple(text_regions),
        containers=containers,
        connectors=tuple(connectors),
        family_proposal=proposal,
        label_font_size=label_font_size,
        container_styles=container_styles,
    )


def _generate_table_matrix_spec(
    rng: random.Random,
    *,
    sample_id: str,
    image_size: AnnotationImageSize,
) -> SyntheticSlideSpec:
    """A regular grid of cell nodes — no connectors, no container."""
    rows = rng.randint(2, 3)
    cols = rng.randint(2, 4)
    node_kind = NodeKind.BOX if rng.random() < 0.6 else NodeKind.ROUNDED_BOX
    theme = _pick_theme(rng)
    _fill, _stroke = _theme_fill_stroke(rng, theme)
    style = SyntheticNodeStyle(fill=_fill, stroke=_stroke, stroke_width=theme.node_stroke_width)

    width, height = float(image_size.width), float(image_size.height)
    margin_x, margin_y = width * 0.1, height * 0.12
    cell_w = (width - 2 * margin_x) / cols
    cell_h = (height - 2 * margin_y) / rows
    gap_x = cell_w * rng.uniform(0.12, 0.22)
    gap_y = cell_h * rng.uniform(0.18, 0.28)
    label_font_size = max(14, int((cell_h - gap_y) * 0.28))
    font = ImageFont.load_default(size=label_font_size)
    labels = rng.sample(_LABEL_VOCAB, rows * cols)

    nodes: list[AnnotationNode] = []
    text_regions: list[AnnotationTextRegion] = []
    node_styles: dict[str, SyntheticNodeStyle] = {}
    index = 0
    for row in range(rows):
        for col in range(cols):
            x0 = margin_x + col * cell_w + gap_x / 2.0
            y0 = margin_y + row * cell_h + gap_y / 2.0
            bbox = AnnotationBBox(x0=x0, y0=y0, x1=x0 + cell_w - gap_x, y1=y0 + cell_h - gap_y)
            node_id = f"node:{sample_id}:{index}"
            text_id = f"text:{sample_id}:{index}"
            label = labels[index]
            nodes.append(
                AnnotationNode(
                    id=node_id, kind=node_kind, bbox=bbox, confidence=1.0, label=label,
                    text_region_ids=(text_id,), source="synthetic_gt", provenance=(f"{GENERATOR_NAME}:node",),
                )
            )
            node_styles[node_id] = style
            text_regions.append(
                AnnotationTextRegion(
                    id=text_id, bbox=_label_bbox(label, font=font, node_bbox=bbox), confidence=1.0,
                    role=TextRegionRole.LABEL, text=label, source="synthetic_gt",
                    provenance=(f"{GENERATOR_NAME}:text",),
                )
            )
            index += 1

    focus = _union_bbox([node.bbox for node in nodes])
    proposal = AnnotationFamilyProposal(
        id=f"family:{sample_id}:0", family=DiagramFamily.TABLE_MATRIX, confidence=1.0, focus_bbox=focus,
        evidence=(f"{GENERATOR_NAME}:layout",), provenance=(f"{GENERATOR_NAME}:family_proposal",),
    )
    return SyntheticSlideSpec(
        sample_id=sample_id, family=DiagramFamily.TABLE_MATRIX, image_size=image_size, nodes=tuple(nodes),
        node_styles=node_styles, text_regions=tuple(text_regions), containers=(), connectors=(),
        family_proposal=proposal, label_font_size=label_font_size, container_styles=(),
    )


def _generate_block_flow_spec(
    rng: random.Random,
    *,
    sample_id: str,
    image_size: AnnotationImageSize,
) -> SyntheticSlideSpec:
    """A shallow tree: a root branches to children (and optionally grandchildren)."""
    node_kind = NodeKind.ROUNDED_BOX if rng.random() < 0.5 else NodeKind.BOX
    theme = _pick_theme(rng)
    _fill, _stroke = _theme_fill_stroke(rng, theme)
    style = SyntheticNodeStyle(fill=_fill, stroke=_stroke, stroke_width=theme.node_stroke_width)

    width, height = float(image_size.width), float(image_size.height)
    child_count = rng.randint(2, 3)
    grand_count = rng.randint(0, 2)
    levels = 3 if grand_count > 0 else 2
    node_w, node_h = width * 0.16, height * 0.16
    margin_y = height * 0.1
    level_gap = (height - 2 * margin_y - node_h) / (levels - 1)
    label_font_size = max(14, int(node_h * 0.28))
    font = ImageFont.load_default(size=label_font_size)
    labels = rng.sample(_LABEL_VOCAB, 1 + child_count + grand_count)

    nodes: list[AnnotationNode] = []
    node_styles: dict[str, SyntheticNodeStyle] = {}
    text_regions: list[AnnotationTextRegion] = []

    def add_node(index: int, center_x: float, center_y: float) -> AnnotationNode:
        bbox = AnnotationBBox(
            x0=center_x - node_w / 2.0, y0=center_y - node_h / 2.0,
            x1=center_x + node_w / 2.0, y1=center_y + node_h / 2.0,
        )
        node_id = f"node:{sample_id}:{index}"
        text_id = f"text:{sample_id}:{index}"
        label = labels[index]
        node = AnnotationNode(
            id=node_id, kind=node_kind, bbox=bbox, confidence=1.0, label=label,
            text_region_ids=(text_id,), source="synthetic_gt", provenance=(f"{GENERATOR_NAME}:node",),
        )
        nodes.append(node)
        node_styles[node_id] = style
        text_regions.append(
            AnnotationTextRegion(
                id=text_id, bbox=_label_bbox(label, font=font, node_bbox=bbox), confidence=1.0,
                role=TextRegionRole.LABEL, text=label, source="synthetic_gt",
                provenance=(f"{GENERATOR_NAME}:text",),
            )
        )
        return node

    root = add_node(0, width / 2.0, margin_y + node_h / 2.0)
    child_y = margin_y + node_h / 2.0 + level_gap
    children = [add_node(1 + i, width * (i + 1) / (child_count + 1), child_y) for i in range(child_count)]
    grandchildren: list[AnnotationNode] = []
    if grand_count > 0:
        base_cx = (children[0].bbox.x0 + children[0].bbox.x1) / 2.0
        grand_y = margin_y + node_h / 2.0 + 2 * level_gap
        for j in range(grand_count):
            gx = base_cx + (j - (grand_count - 1) / 2.0) * node_w * 1.3
            grandchildren.append(add_node(1 + child_count + j, gx, grand_y))

    edges = [(root, child) for child in children] + [(children[0], g) for g in grandchildren]
    connectors: list[SyntheticConnector] = []
    for index, (parent, child) in enumerate(edges):
        parent_center = AnnotationPoint((parent.bbox.x0 + parent.bbox.x1) / 2.0, (parent.bbox.y0 + parent.bbox.y1) / 2.0)
        child_center = AnnotationPoint((child.bbox.x0 + child.bbox.x1) / 2.0, (child.bbox.y0 + child.bbox.y1) / 2.0)
        start_point, start_side = _edge_point_toward(parent.bbox, child_center)
        end_point, end_side = _edge_point_toward(child.bbox, parent_center)
        connector_id = f"connector:{sample_id}:{index}"
        path = (start_point, end_point)
        connectors.append(
            SyntheticConnector(
                candidate=AnnotationConnectorCandidate(
                    id=connector_id, kind=ConnectorKind.ARROW, bbox=_path_bbox(path), confidence=1.0,
                    source_evidence_id=f"evidence:{connector_id}", path_points=path,
                    start_endpoint=AnnotationConnectorEndpoint(
                        point=start_point, owner_id=parent.id, owner_kind=PortOwnerKind.NODE, side=start_side),
                    end_endpoint=AnnotationConnectorEndpoint(
                        point=end_point, owner_id=child.id, owner_kind=PortOwnerKind.NODE, side=end_side),
                    arrowhead_end=True, source="synthetic_gt", provenance=(f"{GENERATOR_NAME}:connector",),
                ),
                start_port=_port_for(parent.id, sample_id, index, "start", side=start_side, point=start_point),
                end_port=_port_for(child.id, sample_id, index, "end", side=end_side, point=end_point),
                stroke=_connector_stroke(theme, style.stroke),
                width=theme.connector_width,
            )
        )

    containers, container_styles = _build_containers(
        rng, sample_id=sample_id, nodes=nodes, image_size=image_size, probability=0.5
    )
    focus = containers[0].bbox if containers else _union_bbox(
        [node.bbox for node in nodes] + [connector.candidate.bbox for connector in connectors]
    )
    proposal = AnnotationFamilyProposal(
        id=f"family:{sample_id}:0", family=DiagramFamily.BLOCK_FLOW, confidence=1.0, focus_bbox=focus,
        evidence=(f"{GENERATOR_NAME}:layout",), provenance=(f"{GENERATOR_NAME}:family_proposal",),
    )
    return SyntheticSlideSpec(
        sample_id=sample_id, family=DiagramFamily.BLOCK_FLOW, image_size=image_size, nodes=tuple(nodes),
        node_styles=node_styles, text_regions=tuple(text_regions), containers=containers, connectors=tuple(connectors),
        family_proposal=proposal, label_font_size=label_font_size, container_styles=container_styles,
    )


def _make_connector(
    *,
    sample_id: str,
    index: int,
    start_node: AnnotationNode,
    end_node: AnnotationNode,
    start_point: AnnotationPoint,
    start_side: PortSide,
    end_point: AnnotationPoint,
    end_side: PortSide,
    theme: RenderTheme,
    node_stroke: tuple[int, int, int],
    path: tuple[AnnotationPoint, ...] | None = None,
    draw_arrowhead: bool = True,
    stroke_override: tuple[int, int, int] | None = None,
    width_override: int | None = None,
) -> SyntheticConnector:
    """Build a directed SyntheticConnector + its ports with theme-aware styling.

    ``stroke_override`` / ``width_override`` force a specific connector look (e.g.
    the solid-black elbow connectors of PowerPoint-style block diagrams).
    """
    resolved_path = path if path is not None else (start_point, end_point)
    connector_id = f"connector:{sample_id}:{index}"
    return SyntheticConnector(
        candidate=AnnotationConnectorCandidate(
            id=connector_id,
            kind=ConnectorKind.ARROW,
            bbox=_path_bbox(resolved_path),
            confidence=1.0,
            source_evidence_id=f"evidence:{connector_id}",
            path_points=resolved_path,
            start_endpoint=AnnotationConnectorEndpoint(
                point=start_point, owner_id=start_node.id, owner_kind=PortOwnerKind.NODE, side=start_side
            ),
            end_endpoint=AnnotationConnectorEndpoint(
                point=end_point, owner_id=end_node.id, owner_kind=PortOwnerKind.NODE, side=end_side
            ),
            arrowhead_end=draw_arrowhead,
            source="synthetic_gt",
            provenance=(f"{GENERATOR_NAME}:connector",),
        ),
        start_port=_port_for(start_node.id, sample_id, index, "start", side=start_side, point=start_point),
        end_port=_port_for(end_node.id, sample_id, index, "end", side=end_side, point=end_point),
        stroke=stroke_override if stroke_override is not None else _connector_stroke(theme, node_stroke),
        width=width_override if width_override is not None else theme.connector_width,
        draw_arrowhead=draw_arrowhead,
    )


def _side_point(box: AnnotationBBox, side: PortSide) -> AnnotationPoint:
    mid_x = (box.x0 + box.x1) / 2.0
    mid_y = (box.y0 + box.y1) / 2.0
    if side is PortSide.RIGHT:
        return AnnotationPoint(box.x1, mid_y)
    if side is PortSide.LEFT:
        return AnnotationPoint(box.x0, mid_y)
    if side is PortSide.TOP:
        return AnnotationPoint(mid_x, box.y0)
    return AnnotationPoint(mid_x, box.y1)


def _elbow_path(
    start: AnnotationPoint, start_side: PortSide, end: AnnotationPoint, end_side: PortSide
) -> tuple[AnnotationPoint, ...]:
    """Right-angle (elbow) route between two side points — the PowerPoint look."""
    horizontal_sides = (PortSide.LEFT, PortSide.RIGHT)
    if start_side in horizontal_sides and end_side in horizontal_sides:
        mid_x = (start.x + end.x) / 2.0
        return (start, AnnotationPoint(mid_x, start.y), AnnotationPoint(mid_x, end.y), end)
    if start_side not in horizontal_sides and end_side not in horizontal_sides:
        mid_y = (start.y + end.y) / 2.0
        return (start, AnnotationPoint(start.x, mid_y), AnnotationPoint(end.x, mid_y), end)
    if start_side in horizontal_sides:
        return (start, AnnotationPoint(end.x, start.y), end)
    return (start, AnnotationPoint(start.x, end.y), end)


def _make_text_region(label: str, *, sample_id: str, index: int, font: ImageFont.ImageFont | ImageFont.FreeTypeFont, node_bbox: AnnotationBBox) -> AnnotationTextRegion:
    return AnnotationTextRegion(
        id=f"text:{sample_id}:{index}",
        bbox=_label_bbox(label, font=font, node_bbox=node_bbox),
        confidence=1.0,
        role=TextRegionRole.LABEL,
        text=label,
        source="synthetic_gt",
        provenance=(f"{GENERATOR_NAME}:text",),
    )


def _container_for(
    members: list[AnnotationNode],
    *,
    sample_id: str,
    index: int,
    image_size: AnnotationImageSize,
    pad_scale: float,
    kind: ContainerKind,
) -> AnnotationContainer:
    width, height = float(image_size.width), float(image_size.height)
    pad = min(width, height) * pad_scale
    union = _union_bbox([node.bbox for node in members])
    return AnnotationContainer(
        id=f"container:{sample_id}:{index}",
        kind=kind,
        bbox=AnnotationBBox(
            x0=max(2.0, union.x0 - pad),
            y0=max(2.0, union.y0 - pad),
            x1=min(width - 2.0, union.x1 + pad),
            y1=min(height - 2.0, union.y1 + pad),
        ),
        confidence=1.0,
        member_node_ids=tuple(node.id for node in members),
        source="synthetic_gt",
        provenance=(f"{GENERATOR_NAME}:container",),
    )


def _build_containers(
    rng: random.Random,
    *,
    sample_id: str,
    nodes: list[AnnotationNode],
    image_size: AnnotationImageSize,
    probability: float,
) -> tuple[tuple[AnnotationContainer, ...], tuple[SyntheticContainerStyle, ...]]:
    """Build 0-2 group containers (panels) around the nodes.

    Real paper figures group nodes in grey/dashed panels — frequently, sometimes
    nested. Emits an outer panel around all nodes and occasionally a nested inner
    panel around a contiguous subset, each with an independent style. ``nodes``
    list order is used for the subset so flows/stacks get a spatially-contiguous
    inner panel; overlapping containers are permitted by the contract.
    """
    if not nodes or rng.random() >= probability:
        return (), ()
    containers: list[AnnotationContainer] = []
    styles: list[SyntheticContainerStyle] = []
    containers.append(
        _container_for(
            list(nodes), sample_id=sample_id, index=0, image_size=image_size,
            pad_scale=rng.uniform(0.03, 0.06), kind=ContainerKind.PANEL,
        )
    )
    styles.append(_pick_container_style(rng))
    if len(nodes) >= 3 and rng.random() < 0.35:
        cut = rng.randint(2, len(nodes) - 1)
        containers.append(
            _container_for(
                list(nodes[:cut]), sample_id=sample_id, index=1, image_size=image_size,
                pad_scale=rng.uniform(0.012, 0.025), kind=ContainerKind.FLOW_CLUSTER,
            )
        )
        styles.append(_pick_container_style(rng))
    return tuple(containers), tuple(styles)


def _generate_graph_spec(
    rng: random.Random,
    *,
    sample_id: str,
    image_size: AnnotationImageSize,
) -> SyntheticSlideSpec:
    """A general node-link graph: ellipse 'state' nodes with arbitrary directed edges.

    Distinct from CYCLE (a single closed ring): edges form a connected graph with
    extra cross-links, mirroring FSMs / dependency graphs common in NLP papers.
    """
    node_count = rng.randint(4, 7)
    theme = _pick_theme(rng)
    fill, stroke = _theme_fill_stroke(rng, theme)
    node_shape = "ellipse" if rng.random() < 0.8 else ("round" if rng.random() < 0.5 else "rect")
    style = SyntheticNodeStyle(fill=fill, stroke=stroke, shape=node_shape, stroke_width=theme.node_stroke_width)

    width, height = float(image_size.width), float(image_size.height)
    center_x, center_y = width / 2.0, height / 2.0
    radius_x = width * rng.uniform(0.28, 0.36)
    radius_y = height * rng.uniform(0.26, 0.34)
    node_half = min(width, height) * rng.uniform(0.055, 0.08)
    start_angle = -math.pi / 2.0 + rng.uniform(-0.3, 0.3)

    label_font_size = max(14, int(node_half * 2.0 * 0.32))
    font = ImageFont.load_default(size=label_font_size)
    labels = rng.sample(_GRAPH_LABEL_VOCAB, node_count)

    nodes: list[AnnotationNode] = []
    text_regions: list[AnnotationTextRegion] = []
    node_styles: dict[str, SyntheticNodeStyle] = {}
    centers: list[AnnotationPoint] = []
    for index in range(node_count):
        angle = start_angle + 2.0 * math.pi * index / node_count
        cx = center_x + radius_x * math.cos(angle) + rng.uniform(-node_half * 0.4, node_half * 0.4)
        cy = center_y + radius_y * math.sin(angle) + rng.uniform(-node_half * 0.4, node_half * 0.4)
        centers.append(AnnotationPoint(cx, cy))
        bbox = AnnotationBBox(x0=cx - node_half, y0=cy - node_half, x1=cx + node_half, y1=cy + node_half)
        node_id = f"node:{sample_id}:{index}"
        nodes.append(
            AnnotationNode(
                id=node_id, kind=NodeKind.BOX, bbox=bbox, confidence=1.0, label=labels[index],
                text_region_ids=(f"text:{sample_id}:{index}",), source="synthetic_gt",
                provenance=(f"{GENERATOR_NAME}:node",),
            )
        )
        node_styles[node_id] = style
        text_regions.append(_make_text_region(labels[index], sample_id=sample_id, index=index, font=font, node_bbox=bbox))

    # Edges: a connected backbone (each node links to a random earlier one) plus a
    # few random extra directed cross-links — a graph, not a ring.
    edges: list[tuple[int, int]] = [(rng.randrange(i), i) for i in range(1, node_count)]
    existing = {frozenset(edge) for edge in edges}
    extra = rng.randint(1, max(1, node_count // 2))
    for _ in range(extra):
        a, b = rng.sample(range(node_count), 2)
        if frozenset((a, b)) not in existing:
            edges.append((a, b))
            existing.add(frozenset((a, b)))

    connectors: list[SyntheticConnector] = []
    for index, (a, b) in enumerate(edges):
        start_point, start_side = _edge_point_toward(nodes[a].bbox, centers[b])
        end_point, end_side = _edge_point_toward(nodes[b].bbox, centers[a])
        connectors.append(
            _make_connector(
                sample_id=sample_id, index=index, start_node=nodes[a], end_node=nodes[b],
                start_point=start_point, start_side=start_side, end_point=end_point, end_side=end_side,
                theme=theme, node_stroke=stroke,
            )
        )

    containers, container_styles = _build_containers(
        rng, sample_id=sample_id, nodes=nodes, image_size=image_size, probability=0.4
    )
    focus = containers[0].bbox if containers else _union_bbox(
        [node.bbox for node in nodes] + [connector.candidate.bbox for connector in connectors]
    )
    proposal = AnnotationFamilyProposal(
        id=f"family:{sample_id}:0", family=DiagramFamily.GRAPH, confidence=1.0, focus_bbox=focus,
        evidence=(f"{GENERATOR_NAME}:layout",), provenance=(f"{GENERATOR_NAME}:family_proposal",),
    )
    return SyntheticSlideSpec(
        sample_id=sample_id, family=DiagramFamily.GRAPH, image_size=image_size, nodes=tuple(nodes),
        node_styles=node_styles, text_regions=tuple(text_regions), containers=containers, connectors=tuple(connectors),
        family_proposal=proposal, label_font_size=label_font_size, container_styles=container_styles,
    )


def _generate_tree_spec(
    rng: random.Random,
    *,
    sample_id: str,
    image_size: AnnotationImageSize,
) -> SyntheticSlideSpec:
    """A hierarchical tree of mostly text-only nodes with thin edges (parse-tree look).

    Critical for real-figure transfer: parse/derivation trees in papers are text
    labels joined by thin lines (no boxes), which a box-only detector misses. Here
    most nodes are LABEL_ANCHOR (text-only, GT bbox = tight text extent), so the
    detector learns to localise text-anchored nodes.
    """
    theme = _pick_theme(rng)
    fill, stroke = _theme_fill_stroke(rng, theme)
    text_only = rng.random() < 0.7
    width, height = float(image_size.width), float(image_size.height)
    max_depth = rng.randint(2, 3)
    budget = rng.randint(5, 10)

    # Build parent/child structure breadth-first under a node budget.
    parents: list[int | None] = [None]
    levels: list[int] = [0]
    children_of: dict[int, list[int]] = {0: []}
    frontier = [0]
    level = 0
    while frontier and len(parents) < budget and level < max_depth:
        next_frontier: list[int] = []
        for parent in frontier:
            if len(parents) >= budget:
                break
            n_children = rng.randint(2, 3) if level == 0 else rng.randint(1, 3)
            for _ in range(n_children):
                if len(parents) >= budget:
                    break
                idx = len(parents)
                parents.append(parent)
                levels.append(level + 1)
                children_of[idx] = []
                children_of[parent].append(idx)
                next_frontier.append(idx)
        frontier = next_frontier
        level += 1

    node_count = len(parents)
    depth = max(levels)
    label_font_size = max(13, int(height * 0.035))
    font = ImageFont.load_default(size=label_font_size)

    # x via post-order leaf centring; y by level.
    leaf_counter = [0]
    leaf_slot: dict[int, float] = {}

    def assign(idx: int) -> float:
        kids = children_of[idx]
        if not kids:
            slot = leaf_counter[0]
            leaf_counter[0] += 1
            leaf_slot[idx] = slot
            return slot
        positions = [assign(child) for child in kids]
        leaf_slot[idx] = sum(positions) / len(positions)
        return leaf_slot[idx]

    assign(0)
    leaf_total = max(1, leaf_counter[0])
    margin_x = width * 0.08
    usable_w = width - 2 * margin_x
    margin_y = height * 0.1
    level_gap = (height - 2 * margin_y) / max(1, depth)

    nodes: list[AnnotationNode] = []
    text_regions: list[AnnotationTextRegion] = []
    node_styles: dict[str, SyntheticNodeStyle] = {}
    centers: list[AnnotationPoint] = []
    for idx in range(node_count):
        label = rng.choice(_TREE_LABEL_VOCAB)
        cx = margin_x + (leaf_slot[idx] + 0.5) * (usable_w / leaf_total)
        cy = margin_y + levels[idx] * level_gap
        left, top, right, bottom = font.getbbox(label)
        text_w, text_h = float(right - left), float(bottom - top)
        if text_only:
            pad_x, pad_y = 3.0, 3.0
            kind, node_shape = NodeKind.LABEL_ANCHOR, "text"
        else:
            pad_x, pad_y = max(8.0, text_w * 0.3), max(6.0, text_h * 0.5)
            kind, node_shape = (NodeKind.ROUNDED_BOX, "round") if rng.random() < 0.5 else (NodeKind.BOX, "rect")
        bbox = AnnotationBBox(
            x0=cx - text_w / 2.0 - pad_x, y0=cy - text_h / 2.0 - pad_y,
            x1=cx + text_w / 2.0 + pad_x, y1=cy + text_h / 2.0 + pad_y,
        )
        centers.append(AnnotationPoint(cx, cy))
        node_id = f"node:{sample_id}:{idx}"
        nodes.append(
            AnnotationNode(
                id=node_id, kind=kind, bbox=bbox, confidence=1.0, label=label,
                text_region_ids=(f"text:{sample_id}:{idx}",), source="synthetic_gt",
                provenance=(f"{GENERATOR_NAME}:node",),
            )
        )
        node_styles[node_id] = SyntheticNodeStyle(
            fill=None if text_only else fill, stroke=None if text_only else stroke,
            shape=node_shape, stroke_width=theme.node_stroke_width,
        )
        text_regions.append(_make_text_region(label, sample_id=sample_id, index=idx, font=font, node_bbox=bbox))

    # Edges parent->child: thin lines, usually no arrowhead (top-down = direction).
    arrowed = rng.random() < 0.3
    connectors: list[SyntheticConnector] = []
    edge_index = 0
    for child in range(1, node_count):
        parent = parents[child]
        assert parent is not None
        start_point, start_side = _edge_point_toward(nodes[parent].bbox, centers[child])
        end_point, end_side = _edge_point_toward(nodes[child].bbox, centers[parent])
        connectors.append(
            _make_connector(
                sample_id=sample_id, index=edge_index, start_node=nodes[parent], end_node=nodes[child],
                start_point=start_point, start_side=start_side, end_point=end_point, end_side=end_side,
                theme=theme, node_stroke=stroke if not text_only else (0, 0, 0), draw_arrowhead=arrowed,
            )
        )
        edge_index += 1

    focus = _union_bbox(
        [node.bbox for node in nodes] + [connector.candidate.bbox for connector in connectors]
    )
    proposal = AnnotationFamilyProposal(
        id=f"family:{sample_id}:0", family=DiagramFamily.TREE, confidence=1.0, focus_bbox=focus,
        evidence=(f"{GENERATOR_NAME}:layout",), provenance=(f"{GENERATOR_NAME}:family_proposal",),
    )
    return SyntheticSlideSpec(
        sample_id=sample_id, family=DiagramFamily.TREE, image_size=image_size, nodes=tuple(nodes),
        node_styles=node_styles, text_regions=tuple(text_regions), containers=(), connectors=tuple(connectors),
        family_proposal=proposal, label_font_size=label_font_size, container_styles=(),
    )


def _generate_layered_stack_spec(
    rng: random.Random,
    *,
    sample_id: str,
    image_size: AnnotationImageSize,
) -> SyntheticSlideSpec:
    """A vertical stack of wide layer bars joined by short arrows (NN architecture)."""
    node_count = rng.randint(3, 6)
    theme = _pick_theme(rng)
    fill, stroke = _theme_fill_stroke(rng, theme)
    node_kind = NodeKind.ROUNDED_BOX if rng.random() < 0.5 else NodeKind.BOX
    node_shape = "round" if node_kind is NodeKind.ROUNDED_BOX else "rect"
    style = SyntheticNodeStyle(fill=fill, stroke=stroke, shape=node_shape, stroke_width=theme.node_stroke_width)

    width, height = float(image_size.width), float(image_size.height)
    bar_w = width * rng.uniform(0.4, 0.58)
    center_x = width / 2.0
    margin_y = height * 0.1
    slot = (height - 2 * margin_y) / node_count
    bar_h = slot * rng.uniform(0.5, 0.66)
    top_down = rng.random() < 0.5

    label_font_size = max(13, int(bar_h * 0.34))
    font = ImageFont.load_default(size=label_font_size)
    labels = rng.sample(_LAYER_LABEL_VOCAB, node_count)

    nodes: list[AnnotationNode] = []
    text_regions: list[AnnotationTextRegion] = []
    node_styles: dict[str, SyntheticNodeStyle] = {}
    for index in range(node_count):
        cy = margin_y + slot * index + slot / 2.0
        bbox = AnnotationBBox(x0=center_x - bar_w / 2.0, y0=cy - bar_h / 2.0, x1=center_x + bar_w / 2.0, y1=cy + bar_h / 2.0)
        node_id = f"node:{sample_id}:{index}"
        nodes.append(
            AnnotationNode(
                id=node_id, kind=node_kind, bbox=bbox, confidence=1.0, label=labels[index],
                text_region_ids=(f"text:{sample_id}:{index}",), source="synthetic_gt",
                provenance=(f"{GENERATOR_NAME}:node",),
            )
        )
        node_styles[node_id] = style
        text_regions.append(_make_text_region(labels[index], sample_id=sample_id, index=index, font=font, node_bbox=bbox))

    order = list(range(node_count)) if top_down else list(reversed(range(node_count)))
    connectors: list[SyntheticConnector] = []
    for index in range(node_count - 1):
        a, b = order[index], order[index + 1]
        start_point, start_side = _edge_point_toward(
            nodes[a].bbox, AnnotationPoint((nodes[b].bbox.x0 + nodes[b].bbox.x1) / 2.0, (nodes[b].bbox.y0 + nodes[b].bbox.y1) / 2.0)
        )
        end_point, end_side = _edge_point_toward(
            nodes[b].bbox, AnnotationPoint((nodes[a].bbox.x0 + nodes[a].bbox.x1) / 2.0, (nodes[a].bbox.y0 + nodes[a].bbox.y1) / 2.0)
        )
        connectors.append(
            _make_connector(
                sample_id=sample_id, index=index, start_node=nodes[a], end_node=nodes[b],
                start_point=start_point, start_side=start_side, end_point=end_point, end_side=end_side,
                theme=theme, node_stroke=stroke,
            )
        )

    containers, container_styles = _build_containers(
        rng, sample_id=sample_id, nodes=nodes, image_size=image_size, probability=0.7
    )
    focus = containers[0].bbox if containers else _union_bbox(
        [node.bbox for node in nodes] + [connector.candidate.bbox for connector in connectors]
    )
    proposal = AnnotationFamilyProposal(
        id=f"family:{sample_id}:0", family=DiagramFamily.LAYERED_STACK, confidence=1.0, focus_bbox=focus,
        evidence=(f"{GENERATOR_NAME}:layout",), provenance=(f"{GENERATOR_NAME}:family_proposal",),
    )
    return SyntheticSlideSpec(
        sample_id=sample_id, family=DiagramFamily.LAYERED_STACK, image_size=image_size, nodes=tuple(nodes),
        node_styles=node_styles, text_regions=tuple(text_regions), containers=containers, connectors=tuple(connectors),
        family_proposal=proposal, label_font_size=label_font_size, container_styles=container_styles,
    )


def _edge_point_toward(bbox: AnnotationBBox, target: AnnotationPoint) -> tuple[AnnotationPoint, PortSide]:
    """Point on ``bbox``'s edge along the ray toward ``target``, and that edge's side."""
    center_x = (bbox.x0 + bbox.x1) / 2.0
    center_y = (bbox.y0 + bbox.y1) / 2.0
    half_w = (bbox.x1 - bbox.x0) / 2.0
    half_h = (bbox.y1 - bbox.y0) / 2.0
    dx = target.x - center_x
    dy = target.y - center_y
    if abs(dx) <= 1e-9 and abs(dy) <= 1e-9:
        # Target coincides with the centre: no ray direction. Default to the top edge.
        return AnnotationPoint(center_x, bbox.y0), PortSide.TOP
    scale_x = half_w / abs(dx) if abs(dx) > 1e-9 else math.inf
    scale_y = half_h / abs(dy) if abs(dy) > 1e-9 else math.inf
    scale = min(scale_x, scale_y)
    point = AnnotationPoint(center_x + dx * scale, center_y + dy * scale)
    return point, edge_side(half_w, half_h, dx, dy)


def edge_side(half_w: float, half_h: float, dx: float, dy: float) -> PortSide:
    """Which edge a ray (dx, dy) from a box centre exits, accounting for aspect ratio.

    Shared by the generator and the connector segmenter's endpoint-side assignment
    so they agree for non-square nodes (a mostly-horizontal direction can still
    exit a wide-but-short box's top edge).
    """
    scale_x = half_w / abs(dx) if abs(dx) > 1e-9 else math.inf
    scale_y = half_h / abs(dy) if abs(dy) > 1e-9 else math.inf
    if scale_x <= scale_y:
        return PortSide.RIGHT if dx >= 0 else PortSide.LEFT
    return PortSide.BOTTOM if dy >= 0 else PortSide.TOP


def _port_for(owner_id: str, sample_id: str, index: int, role: str, *, side: PortSide, point: AnnotationPoint) -> AnnotationPort:
    return AnnotationPort(
        id=f"port:{sample_id}:{index}:{role}",
        owner_id=owner_id,
        owner_kind=PortOwnerKind.NODE,
        side=side,
        point=point,
        confidence=1.0,
        source="synthetic_gt",
        provenance=(f"{GENERATOR_NAME}:port",),
    )


def _orthogonal_path(start: AnnotationPoint, end: AnnotationPoint, *, horizontal: bool) -> tuple[AnnotationPoint, ...]:
    if horizontal and abs(start.y - end.y) > 0.5:
        mid_x = (start.x + end.x) / 2.0
        return (start, AnnotationPoint(mid_x, start.y), AnnotationPoint(mid_x, end.y), end)
    if not horizontal and abs(start.x - end.x) > 0.5:
        mid_y = (start.y + end.y) / 2.0
        return (start, AnnotationPoint(start.x, mid_y), AnnotationPoint(end.x, mid_y), end)
    return (start, end)


def _path_bbox(path: tuple[AnnotationPoint, ...], *, pad: float = 3.0) -> AnnotationBBox:
    xs = [point.x for point in path]
    ys = [point.y for point in path]
    return AnnotationBBox(x0=min(xs) - pad, y0=min(ys) - pad, x1=max(xs) + pad, y1=max(ys) + pad)


def _union_bbox(boxes: list[AnnotationBBox]) -> AnnotationBBox:
    return AnnotationBBox(
        x0=min(box.x0 for box in boxes),
        y0=min(box.y0 for box in boxes),
        x1=max(box.x1 for box in boxes),
        y1=max(box.y1 for box in boxes),
    )


def _label_bbox(label: str, *, font: ImageFont.ImageFont | ImageFont.FreeTypeFont, node_bbox: AnnotationBBox) -> AnnotationBBox:
    left, top, right, bottom = font.getbbox(label)
    text_width, text_height = float(right - left), float(bottom - top)
    center_x = (node_bbox.x0 + node_bbox.x1) / 2.0
    center_y = (node_bbox.y0 + node_bbox.y1) / 2.0
    return AnnotationBBox(
        x0=center_x - text_width / 2.0,
        y0=center_y - text_height / 2.0,
        x1=center_x + text_width / 2.0,
        y1=center_y + text_height / 2.0,
    )


def render_spec_image(spec: SyntheticSlideSpec) -> Image.Image:
    image = Image.new("RGB", (spec.image_size.width, spec.image_size.height), (255, 255, 255))
    draw = ImageDraw.Draw(image)

    # Containers first (behind nodes/connectors). Outer panels precede nested ones.
    for index, container in enumerate(spec.containers):
        box = container.bbox
        style = spec.container_styles[index] if index < len(spec.container_styles) else None
        fill = style.fill if style is not None else CONTAINER_FILL
        outline = style.outline if style is not None else CONTAINER_OUTLINE
        outline_width = style.outline_width if style is not None else CONTAINER_OUTLINE_WIDTH
        if style is not None and style.dashed:
            draw.rectangle((box.x0, box.y0, box.x1, box.y1), fill=fill)
            _draw_dashed_rectangle(draw, box, color=outline, width=outline_width)
        elif style is not None and style.shape == "rect":
            draw.rectangle((box.x0, box.y0, box.x1, box.y1), fill=fill, outline=outline, width=outline_width)
        else:
            draw.rounded_rectangle(
                (box.x0, box.y0, box.x1, box.y1), radius=12, fill=fill, outline=outline, width=outline_width
            )

    for connector in spec.connectors:
        points = [(point.x, point.y) for point in connector.candidate.path_points]
        draw.line(points, fill=connector.stroke, width=connector.width)
        if connector.draw_arrowhead and connector.candidate.arrowhead_end:
            _draw_arrowhead(draw, points[-2], points[-1], color=connector.stroke, size=max(8.0, connector.width * 4.0))

    for node in spec.nodes:
        style = spec.node_styles[node.id]
        box = (node.bbox.x0, node.bbox.y0, node.bbox.x1, node.bbox.y1)
        shape = _resolve_shape(style, node.kind)
        if shape == "text":
            continue  # text-only node: the label (drawn below) is the whole node
        if shape == "ellipse":
            draw.ellipse(box, fill=style.fill, outline=style.stroke, width=style.stroke_width)
        elif shape == "round":
            draw.rounded_rectangle(
                box, radius=int(node.bbox.height * 0.2), fill=style.fill, outline=style.stroke, width=style.stroke_width
            )
        else:
            draw.rectangle(box, fill=style.fill, outline=style.stroke, width=style.stroke_width)

    # Non-GT decorations (braces, margin lines) on top — hard negatives, not connectors.
    for deco in spec.decorations:
        draw.line([(float(px), float(py)) for px, py in deco.points], fill=deco.stroke, width=deco.width)

    font = ImageFont.load_default(size=spec.label_font_size)
    for text_region in spec.text_regions:
        if text_region.text is None:
            continue
        left, top, _, _ = font.getbbox(text_region.text)
        draw.text(
            (text_region.bbox.x0 - left, text_region.bbox.y0 - top),
            text_region.text,
            fill=(15, 23, 42),
            font=font,
        )

    return image


def _resolve_shape(style: SyntheticNodeStyle, kind: NodeKind) -> str:
    """Resolve a node's render shape, mapping the ``auto`` default from NodeKind."""
    if style.shape != "auto":
        return style.shape
    if kind is NodeKind.ROUNDED_BOX:
        return "round"
    if kind is NodeKind.LABEL_ANCHOR:
        return "text"
    return "rect"


def _draw_dashed_rectangle(
    draw: ImageDraw.ImageDraw,
    box: AnnotationBBox,
    *,
    color: tuple[int, int, int],
    width: int,
    dash: int = 12,
    gap: int = 8,
) -> None:
    """Draw a dashed rectangle outline (PIL has no native dashed stroke)."""
    x0, y0, x1, y1 = box.x0, box.y0, box.x1, box.y1
    step = dash + gap

    def dashed_segment(ax: float, ay: float, bx: float, by: float) -> None:
        length = math.hypot(bx - ax, by - ay)
        if length <= 0:
            return
        ux, uy = (bx - ax) / length, (by - ay) / length
        pos = 0.0
        while pos < length:
            end = min(pos + dash, length)
            draw.line(
                [(ax + ux * pos, ay + uy * pos), (ax + ux * end, ay + uy * end)], fill=color, width=width
            )
            pos += step

    dashed_segment(x0, y0, x1, y0)
    dashed_segment(x1, y0, x1, y1)
    dashed_segment(x1, y1, x0, y1)
    dashed_segment(x0, y1, x0, y0)


def _draw_arrowhead(
    draw: ImageDraw.ImageDraw,
    tail: tuple[float, float],
    tip: tuple[float, float],
    *,
    color: tuple[int, int, int],
    size: float = 10.0,
) -> None:
    dx, dy = tip[0] - tail[0], tip[1] - tail[1]
    length = max((dx * dx + dy * dy) ** 0.5, 1e-6)
    ux, uy = dx / length, dy / length
    px, py = -uy, ux
    base_x, base_y = tip[0] - ux * size, tip[1] - uy * size
    half = size * 0.5
    draw.polygon(
        [tip, (base_x + px * half, base_y + py * half), (base_x - px * half, base_y - py * half)],
        fill=color,
    )


def write_spec_pptx(spec: SyntheticSlideSpec, path: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Pt

    presentation = Presentation()
    presentation.slide_width = Emu(spec.image_size.width * EMU_PER_PIXEL)
    presentation.slide_height = Emu(spec.image_size.height * EMU_PER_PIXEL)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    def emu_box(bbox: AnnotationBBox) -> tuple[Emu, Emu, Emu, Emu]:
        return (
            Emu(int(bbox.x0 * EMU_PER_PIXEL)),
            Emu(int(bbox.y0 * EMU_PER_PIXEL)),
            Emu(int(bbox.width * EMU_PER_PIXEL)),
            Emu(int(bbox.height * EMU_PER_PIXEL)),
        )

    for c_index, container in enumerate(spec.containers):
        style = spec.container_styles[c_index] if c_index < len(spec.container_styles) else None
        fill = style.fill if style is not None else CONTAINER_FILL
        outline = style.outline if style is not None else CONTAINER_OUTLINE
        outline_width = style.outline_width if style is not None else CONTAINER_OUTLINE_WIDTH
        mso_shape = (
            MSO_SHAPE.RECTANGLE
            if style is not None and (style.dashed or style.shape == "rect")
            else MSO_SHAPE.ROUNDED_RECTANGLE
        )
        shape = slide.shapes.add_shape(mso_shape, *emu_box(container.bbox))
        shape.name = container.id
        shape.shadow.inherit = False
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill)
        shape.line.color.rgb = RGBColor(*outline)
        shape.line.width = Pt(float(outline_width) * 0.75)  # px stroke -> pt at 96 dpi
        if style is not None and style.dashed:
            line_element = shape.line._get_or_add_ln()
            dash = line_element.makeelement(qn("a:prstDash"), {"val": "dash"})
            line_element.append(dash)

    for node in spec.nodes:
        style = spec.node_styles[node.id]
        shape_name = _resolve_shape(style, node.kind)
        if shape_name == "text":
            # Text-only node (e.g. parse-tree label): a borderless text box whose
            # extent matches the GT node bbox so the pptx mirrors the render.
            shape = slide.shapes.add_textbox(*emu_box(node.bbox))
            shape.name = node.id
            shape.fill.background()
            shape.line.fill.background()
        else:
            if shape_name == "ellipse":
                mso_shape = MSO_SHAPE.OVAL
            elif shape_name == "round":
                mso_shape = MSO_SHAPE.ROUNDED_RECTANGLE
            else:
                mso_shape = MSO_SHAPE.RECTANGLE
            shape = slide.shapes.add_shape(mso_shape, *emu_box(node.bbox))
            shape.name = node.id
            shape.shadow.inherit = False
            if style.fill is not None:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*style.fill)
            else:
                shape.fill.background()
            if style.stroke is not None:
                shape.line.color.rgb = RGBColor(*style.stroke)
                shape.line.width = Pt(float(style.stroke_width) * 0.75)
            else:
                shape.line.fill.background()
        if node.label:
            text_frame = shape.text_frame
            text_frame.text = node.label
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Keep the label on one line centred on the node. Tight text-only nodes
            # (parse-tree labels) otherwise wrap in LibreOffice — its font metrics
            # differ from PIL's — pushing glyphs below the GT bbox.
            text_frame.word_wrap = False
            text_frame.margin_left = Emu(0)
            text_frame.margin_right = Emu(0)
            text_frame.margin_top = Emu(0)
            text_frame.margin_bottom = Emu(0)
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            run = paragraph.runs[0]
            # 1 pt = 4/3 px at 96 dpi; keep the pptx text height aligned
            # with the rendered label so GT bboxes stay meaningful.
            run.font.size = Pt(max(8, round(spec.label_font_size * 0.75)))
            run.font.color.rgb = RGBColor(15, 23, 42)

    for connector in spec.connectors:
        points = connector.candidate.effective_path_points()
        for start, end in zip(points, points[1:]):
            shape = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(int(start.x * EMU_PER_PIXEL)),
                Emu(int(start.y * EMU_PER_PIXEL)),
                Emu(int(end.x * EMU_PER_PIXEL)),
                Emu(int(end.y * EMU_PER_PIXEL)),
            )
            shape.name = connector.candidate.id
            shape.line.color.rgb = RGBColor(*connector.stroke)
            shape.line.width = Pt(max(1.0, float(connector.width)) * 0.75)
            if connector.draw_arrowhead and connector.candidate.arrowhead_end and end is points[-1]:
                line_element = shape.line._get_or_add_ln()
                tail = line_element.makeelement(qn("a:tailEnd"), {"type": "triangle"})
                line_element.append(tail)

    # Non-GT decorations (braces, margin lines): rendered as plain straight segments
    # with no arrowhead, never added to the ground-truth connector set.
    for d_index, deco in enumerate(spec.decorations):
        for start, end in zip(deco.points, deco.points[1:]):
            shape = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Emu(int(start[0] * EMU_PER_PIXEL)),
                Emu(int(start[1] * EMU_PER_PIXEL)),
                Emu(int(end[0] * EMU_PER_PIXEL)),
                Emu(int(end[1] * EMU_PER_PIXEL)),
            )
            shape.name = f"decoration:{d_index}"
            shape.line.color.rgb = RGBColor(*deco.stroke)
            shape.line.width = Pt(max(1.0, float(deco.width)) * 0.75)

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def find_soffice() -> str | None:
    """Locate a LibreOffice binary for pptx -> png rendering, if installed."""
    import shutil

    binary = shutil.which("soffice") or shutil.which("soffice.exe")
    if binary:
        return binary
    candidates = (
        Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"),
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
    )
    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    return None


def render_pptx_batch_via_soffice(
    pptx_paths: list[Path],
    *,
    output_dir: Path,
    image_size: AnnotationImageSize,
    soffice_binary: str | None = None,
) -> None:
    """Convert pptx files to png in one soffice invocation (per-call startup is slow).

    LibreOffice exports at 96 dpi, which matches the 1 px = 9525 EMU slide
    contract, so output size normally equals ``image_size`` already; resize
    only as a safety net.
    """
    import subprocess

    binary = soffice_binary or find_soffice()
    if binary is None:
        raise RuntimeError("LibreOffice (soffice) not found; install it or use the pil renderer")
    if not pptx_paths:
        return
    output_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [binary, "--headless", "--convert-to", "png", "--outdir", str(output_dir), *map(str, pptx_paths)],
        check=True,
        capture_output=True,
        timeout=120 + 10 * len(pptx_paths),
    )
    expected = (image_size.width, image_size.height)
    for pptx_path in pptx_paths:
        png_path = output_dir / f"{pptx_path.stem}.png"
        if not png_path.exists():
            raise RuntimeError(f"soffice did not produce {png_path}")
        with Image.open(png_path) as image:
            if image.size != expected:
                image.convert("RGB").resize(expected, Image.LANCZOS).save(png_path)


def validate_spec_contract(spec: SyntheticSlideSpec) -> None:
    """Guarantee generated ground truth satisfies the v3 SlideIR contract."""
    document = spec.to_annotation_document()
    slide_ir = _GT_ADAPTER.to_slide_ir(document)
    validate_slide_ir(slide_ir)
