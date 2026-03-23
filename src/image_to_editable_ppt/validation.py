from __future__ import annotations

from dataclasses import asdict, dataclass
import json
import math
from pathlib import Path

import numpy as np
from PIL import Image, ImageColor, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt

from .config import PipelineConfig
from .ir import BBox, Point
from .pipeline import ConversionResult, convert_image
from .preprocess import load_image, preprocess_image
from .style import dilate_mask
from .svg_exporter import SVG_NS, format_number, to_svg_color

XML_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


@dataclass(slots=True, frozen=True)
class ValidationShape:
    kind: str
    bbox: BBox
    stroke_color: tuple[int, int, int]
    stroke_width: float
    fill_color: tuple[int, int, int] | None = None
    points: tuple[Point, ...] = ()
    corner_radius: float = 0.0
    text: str = ""
    arrow_start: bool = False
    arrow_end: bool = False


@dataclass(slots=True, frozen=True)
class ValidationMetrics:
    rendered_shape_count: int
    precision: float
    recall: float
    f1: float
    coverage_ratio: float
    blank_output_penalty: float
    structure_score: float
    input_edge_pixels: int
    output_edge_pixels: int
    overlap_pixels: int


@dataclass(slots=True, frozen=True)
class ValidationArtifacts:
    iteration_dir: Path
    output_pptx: Path
    output_svg: Path
    rendered_png: Path
    overlay_png: Path
    edge_diff_png: Path
    metrics_json: Path
    elements_json: Path
    rejections_json: Path


@dataclass(slots=True, frozen=True)
class ValidationRun:
    conversion: ConversionResult
    shapes: list[ValidationShape]
    metrics: ValidationMetrics
    artifacts: ValidationArtifacts


def run_validation_iteration(
    input_path: str | Path,
    iteration_dir: str | Path,
    *,
    config: PipelineConfig | None = None,
    enable_ocr: bool = False,
) -> ValidationRun:
    active_config = config or PipelineConfig()
    iteration_path = Path(iteration_dir)
    iteration_path.mkdir(parents=True, exist_ok=True)
    elements_path = iteration_path / "elements.json"
    output_pptx = iteration_path / "output.pptx"
    input_image = load_image(input_path)
    conversion = convert_image(
        input_path,
        output_pptx,
        config=active_config,
        enable_ocr=enable_ocr,
        debug_elements_path=elements_path,
    )
    shapes = load_pptx_shapes(output_pptx, image_size=input_image.size, config=active_config)
    output_svg = iteration_path / "output.svg"
    export_validation_svg(shapes, input_image.size, output_svg)
    rendered = render_shapes(shapes, input_image.size)
    rendered_png = iteration_path / "output.png"
    rendered.save(rendered_png)
    overlay = render_overlay(input_image, rendered)
    overlay_png = iteration_path / "overlay.png"
    overlay.save(overlay_png)
    edge_diff, metrics = compare_input_to_render(input_image, rendered, active_config, len(shapes))
    edge_diff_png = iteration_path / "edge-diff.png"
    edge_diff.save(edge_diff_png)
    metrics_json = iteration_path / "comparison.json"
    with metrics_json.open("w", encoding="utf-8") as handle:
        json.dump(asdict(metrics), handle, indent=2)
    return ValidationRun(
        conversion=conversion,
        shapes=shapes,
        metrics=metrics,
        artifacts=ValidationArtifacts(
            iteration_dir=iteration_path,
            output_pptx=output_pptx,
            output_svg=output_svg,
            rendered_png=rendered_png,
            overlay_png=overlay_png,
            edge_diff_png=edge_diff_png,
            metrics_json=metrics_json,
            elements_json=elements_path,
            rejections_json=elements_path.with_name(f"{elements_path.stem}.rejections{elements_path.suffix}"),
        ),
    )


def load_pptx_shapes(
    pptx_path: str | Path,
    *,
    image_size: tuple[int, int],
    config: PipelineConfig,
) -> list[ValidationShape]:
    presentation = Presentation(str(pptx_path))
    slide = presentation.slides[0]
    transform = export_transform(
        image_size=image_size,
        slide_width=int(presentation.slide_width),
        slide_height=int(presentation.slide_height),
        config=config,
    )
    shapes: list[ValidationShape] = []
    for shape in slide.shapes:
        parsed = parse_shape(shape, transform)
        if parsed is None:
            continue
        shapes.append(parsed)
    return shapes


def export_validation_svg(
    shapes: list[ValidationShape],
    image_size: tuple[int, int],
    output_path: str | Path,
) -> None:
    import xml.etree.ElementTree as ET

    width, height = image_size
    ET.register_namespace("", SVG_NS)
    svg = ET.Element(
        f"{{{SVG_NS}}}svg",
        attrib={
            "width": str(width),
            "height": str(height),
            "viewBox": f"0 0 {width} {height}",
            "version": "1.1",
        },
    )
    defs = ET.SubElement(svg, f"{{{SVG_NS}}}defs")
    for marker_id in ("arrow-tip", "arrow-tail"):
        marker = ET.SubElement(
            defs,
            f"{{{SVG_NS}}}marker",
            attrib={
                "id": marker_id,
                "markerWidth": "10",
                "markerHeight": "7",
                "refX": "9" if marker_id == "arrow-tip" else "1",
                "refY": "3.5",
                "orient": "auto",
                "markerUnits": "strokeWidth",
            },
        )
        polygon_points = "0 0, 10 3.5, 0 7" if marker_id == "arrow-tip" else "10 0, 0 3.5, 10 7"
        ET.SubElement(
            marker,
            f"{{{SVG_NS}}}polygon",
            attrib={"points": polygon_points, "fill": "#000000"},
        )
    ET.SubElement(
        svg,
        f"{{{SVG_NS}}}rect",
        attrib={"x": "0", "y": "0", "width": str(width), "height": str(height), "fill": "#ffffff"},
    )
    for shape in shapes:
        append_svg_shape(svg, shape)
    ET.ElementTree(svg).write(Path(output_path), encoding="utf-8", xml_declaration=True)


def render_shapes(
    shapes: list[ValidationShape],
    image_size: tuple[int, int],
    *,
    include_text: bool = True,
) -> Image.Image:
    image = Image.new("RGB", image_size, "white")
    draw = ImageDraw.Draw(image)
    font = ImageFont.load_default()
    for shape in shapes:
        if shape.kind == "rect":
            bbox = bbox_tuple(shape.bbox)
            if shape.corner_radius > 0:
                draw.rounded_rectangle(
                    bbox,
                    radius=int(round(shape.corner_radius)),
                    outline=shape.stroke_color,
                    fill=shape.fill_color,
                    width=max(1, int(round(shape.stroke_width))),
                )
            else:
                draw.rectangle(
                    bbox,
                    outline=shape.stroke_color,
                    fill=shape.fill_color,
                    width=max(1, int(round(shape.stroke_width))),
                )
            continue
        if shape.kind == "line":
            if len(shape.points) < 2:
                continue
            draw.line(flatten_points(shape.points), fill=shape.stroke_color, width=max(1, int(round(shape.stroke_width))))
            if shape.arrow_start:
                draw_arrowhead(draw, shape.points[1], shape.points[0], shape.stroke_color, shape.stroke_width)
            if shape.arrow_end:
                draw_arrowhead(draw, shape.points[-2], shape.points[-1], shape.stroke_color, shape.stroke_width)
            continue
        if shape.kind == "polyline":
            if len(shape.points) < 2:
                continue
            draw.line(flatten_points(shape.points), fill=shape.stroke_color, width=max(1, int(round(shape.stroke_width))))
            if shape.arrow_start and len(shape.points) >= 2:
                draw_arrowhead(draw, shape.points[1], shape.points[0], shape.stroke_color, shape.stroke_width)
            if shape.arrow_end and len(shape.points) >= 2:
                draw_arrowhead(draw, shape.points[-2], shape.points[-1], shape.stroke_color, shape.stroke_width)
            continue
        if shape.kind == "polygon":
            if len(shape.points) < 3:
                continue
            draw.polygon(flatten_points(shape.points), fill=shape.fill_color, outline=shape.stroke_color)
            continue
        if shape.kind == "text" and include_text and shape.text:
            draw.multiline_text(
                (shape.bbox.x0 + 2.0, shape.bbox.y0 + 2.0),
                shape.text,
                fill=shape.stroke_color,
                font=font,
                spacing=2,
            )
    return image


def render_overlay(input_image: Image.Image, rendered_image: Image.Image) -> Image.Image:
    base = Image.blend(input_image.convert("RGB"), Image.new("RGB", input_image.size, "white"), alpha=0.35)
    rendered = np.asarray(rendered_image.convert("RGB"), dtype=np.uint8)
    overlay = np.asarray(base, dtype=np.uint8).copy()
    structure = np.any(rendered < 235, axis=2)
    overlay[structure] = np.asarray((220, 36, 36), dtype=np.uint8)
    return Image.fromarray(overlay)


def compare_input_to_render(
    input_image: Image.Image,
    rendered_image: Image.Image,
    config: PipelineConfig,
    rendered_shape_count: int,
) -> tuple[Image.Image, ValidationMetrics]:
    input_processed = preprocess_image(
        input_image,
        foreground_threshold=config.foreground_threshold,
        min_component_area=config.min_component_area,
        min_stroke_length=config.min_stroke_length,
        min_box_size=config.min_box_size,
        min_relative_line_length=config.min_relative_line_length,
        min_relative_box_size=config.min_relative_box_size,
        adaptive_background=config.adaptive_background,
        background_blur_divisor=config.background_blur_divisor,
    )
    render_processed = preprocess_image(
        rendered_image,
        foreground_threshold=max(14.0, config.foreground_threshold * 0.45),
        min_component_area=config.min_component_area,
        min_stroke_length=config.min_stroke_length,
        min_box_size=config.min_box_size,
        min_relative_line_length=config.min_relative_line_length,
        min_relative_box_size=config.min_relative_box_size,
        adaptive_background=False,
        background_blur_divisor=config.background_blur_divisor,
    )
    input_edges = input_processed.boundary_mask
    output_edges = render_processed.boundary_mask
    input_dilated = dilate_mask(input_edges, radius=2)
    output_dilated = dilate_mask(output_edges, radius=2)
    overlap = output_edges & input_dilated
    precision = float(overlap.sum() / max(1, int(output_edges.sum())))
    recall_overlap = input_edges & output_dilated
    recall = float(recall_overlap.sum() / max(1, int(input_edges.sum())))
    f1 = 0.0 if precision + recall == 0.0 else float(2.0 * precision * recall / (precision + recall))
    coverage_ratio = float(output_edges.sum() / max(1, int(input_edges.sum())))
    blank_output_penalty = max(0.0, (0.02 - coverage_ratio) / 0.02)
    structure_score = 20.0 * recall + 8.0 * precision - 30.0 * blank_output_penalty
    diff = edge_diff_image(input_edges, output_edges)
    return diff, ValidationMetrics(
        rendered_shape_count=rendered_shape_count,
        precision=precision,
        recall=recall,
        f1=f1,
        coverage_ratio=coverage_ratio,
        blank_output_penalty=blank_output_penalty,
        structure_score=structure_score,
        input_edge_pixels=int(input_edges.sum()),
        output_edge_pixels=int(output_edges.sum()),
        overlap_pixels=int(overlap.sum()),
    )


def edge_diff_image(input_edges: np.ndarray, output_edges: np.ndarray) -> Image.Image:
    canvas = np.full((*input_edges.shape, 3), 255, dtype=np.uint8)
    both = input_edges & output_edges
    input_only = input_edges & ~output_edges
    output_only = output_edges & ~input_edges
    canvas[input_only] = np.asarray((216, 46, 46), dtype=np.uint8)
    canvas[output_only] = np.asarray((42, 132, 42), dtype=np.uint8)
    canvas[both] = np.asarray((32, 32, 32), dtype=np.uint8)
    return Image.fromarray(canvas)


def parse_shape(shape, transform: tuple[float, int, int]) -> ValidationShape | None:
    scale, offset_x, offset_y = transform
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        prst = shape_preset(shape)
        if prst not in {"rect", "roundRect"}:
            return None
        bbox = map_bbox(shape.left, shape.top, shape.width, shape.height, scale, offset_x, offset_y)
        stroke_color = line_color(shape._element)
        stroke_width = max(1.0, emu_to_px(line_width(shape._element), scale))
        fill_color = fill_color_from_element(shape._element)
        corner_ratio = 0.12 if prst == "roundRect" else 0.0
        return ValidationShape(
            kind="rect",
            bbox=bbox,
            stroke_color=stroke_color,
            stroke_width=stroke_width,
            fill_color=fill_color,
            corner_radius=min(bbox.width, bbox.height) * corner_ratio,
        )
    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        bbox = map_bbox(shape.left, shape.top, shape.width, shape.height, scale, offset_x, offset_y)
        return ValidationShape(
            kind="text",
            bbox=bbox,
            stroke_color=(0, 0, 0),
            stroke_width=0.0,
            text=shape.text,
        )
    if shape.shape_type == MSO_SHAPE_TYPE.LINE:
        points = connector_points(shape._element, scale, offset_x, offset_y)
        if len(points) != 2:
            return None
        bbox = BBox(
            min(points[0].x, points[1].x),
            min(points[0].y, points[1].y),
            max(points[0].x, points[1].x),
            max(points[0].y, points[1].y),
        )
        return ValidationShape(
            kind="line",
            bbox=bbox,
            stroke_color=line_color(shape._element),
            stroke_width=max(1.0, emu_to_px(line_width(shape._element), scale)),
            points=points,
            arrow_start=has_arrowhead(shape._element, at="head"),
            arrow_end=has_arrowhead(shape._element, at="tail"),
        )
    if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
        points, closed = freeform_points(shape._element, scale, offset_x, offset_y)
        if len(points) < 2:
            return None
        bbox = BBox(
            min(point.x for point in points),
            min(point.y for point in points),
            max(point.x for point in points),
            max(point.y for point in points),
        )
        kind = "polygon" if closed and fill_color_from_element(shape._element) is not None else "polyline"
        return ValidationShape(
            kind=kind,
            bbox=bbox,
            stroke_color=line_color(shape._element),
            stroke_width=max(1.0, emu_to_px(line_width(shape._element), scale)),
            fill_color=fill_color_from_element(shape._element),
            points=points,
        )
    return None


def export_transform(
    *,
    image_size: tuple[int, int],
    slide_width: int,
    slide_height: int,
    config: PipelineConfig,
) -> tuple[float, int, int]:
    padding = int(Pt(config.slide_padding_pt))
    available_width = slide_width - padding * 2
    available_height = slide_height - padding * 2
    image_width, image_height = image_size
    scale = min(available_width / image_width, available_height / image_height)
    offset_x = int((slide_width - image_width * scale) / 2.0)
    offset_y = int((slide_height - image_height * scale) / 2.0)
    return scale, offset_x, offset_y


def shape_preset(shape) -> str | None:
    nodes = shape._element.xpath("./p:spPr/a:prstGeom")
    if not nodes:
        return None
    return nodes[0].get("prst")


def connector_points(element, scale: float, offset_x: int, offset_y: int) -> tuple[Point, Point]:
    xfrm = element.xpath("./p:spPr/a:xfrm")
    if not xfrm:
        return ()
    node = xfrm[0]
    off = node.xpath("./a:off")[0]
    ext = node.xpath("./a:ext")[0]
    x0 = int(off.get("x", "0"))
    y0 = int(off.get("y", "0"))
    cx = int(ext.get("cx", "0"))
    cy = int(ext.get("cy", "0"))
    flip_h = node.get("flipH") == "1"
    flip_v = node.get("flipV") == "1"
    start_x = x0 + (cx if flip_h else 0)
    end_x = x0 + (0 if flip_h else cx)
    start_y = y0 + (cy if flip_v else 0)
    end_y = y0 + (0 if flip_v else cy)
    return (
        Point(emu_to_px(start_x - offset_x, scale), emu_to_px(start_y - offset_y, scale)),
        Point(emu_to_px(end_x - offset_x, scale), emu_to_px(end_y - offset_y, scale)),
    )


def freeform_points(
    element,
    scale: float,
    offset_x: int,
    offset_y: int,
) -> tuple[tuple[Point, ...], bool]:
    xfrm = element.xpath("./p:spPr/a:xfrm")
    path_nodes = element.xpath("./p:spPr/a:custGeom/a:pathLst/a:path")
    if not xfrm or not path_nodes:
        return (), False
    node = xfrm[0]
    off = node.xpath("./a:off")[0]
    ext = node.xpath("./a:ext")[0]
    offset_emu_x = int(off.get("x", "0"))
    offset_emu_y = int(off.get("y", "0"))
    ext_x = max(1, int(ext.get("cx", "1")))
    ext_y = max(1, int(ext.get("cy", "1")))
    flip_h = node.get("flipH") == "1"
    flip_v = node.get("flipV") == "1"
    path = path_nodes[0]
    path_width = max(1, int(path.get("w", str(ext_x))))
    path_height = max(1, int(path.get("h", str(ext_y))))
    points: list[Point] = []
    closed = False
    for child in path:
        local_name = child.tag.rsplit("}", 1)[-1]
        if local_name == "close":
            closed = True
            continue
        pt_nodes = child.xpath("./a:pt")
        if not pt_nodes:
            continue
        pt = pt_nodes[0]
        local_x = int(pt.get("x", "0"))
        local_y = int(pt.get("y", "0"))
        if flip_h:
            local_x = path_width - local_x
        if flip_v:
            local_y = path_height - local_y
        emu_x = offset_emu_x + int(round(local_x * ext_x / path_width))
        emu_y = offset_emu_y + int(round(local_y * ext_y / path_height))
        points.append(
            Point(
                emu_to_px(emu_x - offset_x, scale),
                emu_to_px(emu_y - offset_y, scale),
            )
        )
    if closed and points and points[0] != points[-1]:
        points.append(points[0])
    return tuple(points), closed


def map_bbox(left: Emu, top: Emu, width: Emu, height: Emu, scale: float, offset_x: int, offset_y: int) -> BBox:
    x0 = emu_to_px(int(left) - offset_x, scale)
    y0 = emu_to_px(int(top) - offset_y, scale)
    x1 = emu_to_px(int(left + width) - offset_x, scale)
    y1 = emu_to_px(int(top + height) - offset_y, scale)
    return BBox(x0, y0, x1, y1)


def emu_to_px(value: int | float, scale: float) -> float:
    return float(value) / scale


def line_width(element) -> int:
    nodes = element.xpath("./p:spPr/a:ln")
    if not nodes:
        return 0
    return int(nodes[0].get("w", "0"))


def line_color(element) -> tuple[int, int, int]:
    nodes = element.xpath("./p:spPr/a:ln/a:solidFill/a:srgbClr")
    if not nodes:
        return (0, 0, 0)
    return hex_to_rgb(nodes[0].get("val", "000000"))


def fill_color_from_element(element) -> tuple[int, int, int] | None:
    no_fill = element.xpath("./p:spPr/a:noFill")
    if no_fill:
        return None
    fill = element.xpath("./p:spPr/a:solidFill/a:srgbClr")
    if not fill:
        return None
    return hex_to_rgb(fill[0].get("val", "FFFFFF"))


def has_arrowhead(element, *, at: str) -> bool:
    tag = "headEnd" if at == "head" else "tailEnd"
    return bool(element.xpath(f"./p:spPr/a:ln/a:{tag}"))


def hex_to_rgb(value: str) -> tuple[int, int, int]:
    return ImageColor.getrgb(f"#{value}")


def append_svg_shape(parent, shape: ValidationShape) -> None:
    import xml.etree.ElementTree as ET

    if shape.kind == "rect":
        attrib = {
            "x": format_number(shape.bbox.x0),
            "y": format_number(shape.bbox.y0),
            "width": format_number(shape.bbox.width),
            "height": format_number(shape.bbox.height),
            "fill": to_svg_color(shape.fill_color),
            "stroke": to_svg_color(shape.stroke_color),
            "stroke-width": format_number(max(1.0, shape.stroke_width)),
        }
        if shape.corner_radius > 0:
            attrib["rx"] = format_number(shape.corner_radius)
            attrib["ry"] = format_number(shape.corner_radius)
        ET.SubElement(parent, f"{{{SVG_NS}}}rect", attrib=attrib)
        return
    if shape.kind == "text":
        if not shape.text:
            return
        ET.SubElement(
            parent,
            f"{{{SVG_NS}}}text",
            attrib={
                "x": format_number(shape.bbox.x0 + 2.0),
                "y": format_number(shape.bbox.y0 + max(10.0, shape.bbox.height * 0.6)),
                "fill": to_svg_color(shape.stroke_color),
                "font-size": format_number(max(10.0, shape.bbox.height * 0.7)),
                "font-family": "sans-serif",
            },
        ).text = shape.text
        return
    if shape.kind == "polygon":
        ET.SubElement(
            parent,
            f"{{{SVG_NS}}}polygon",
            attrib={
                "points": point_string(shape.points),
                "fill": to_svg_color(shape.fill_color),
                "stroke": to_svg_color(shape.stroke_color),
                "stroke-width": format_number(max(1.0, shape.stroke_width)),
            },
        )
        return
    tag = "line" if shape.kind == "line" and len(shape.points) == 2 else "polyline"
    if tag == "line":
        start, end = shape.points
        attrib = {
            "x1": format_number(start.x),
            "y1": format_number(start.y),
            "x2": format_number(end.x),
            "y2": format_number(end.y),
            "fill": "none",
            "stroke": to_svg_color(shape.stroke_color),
            "stroke-width": format_number(max(1.0, shape.stroke_width)),
            "stroke-linecap": "round",
        }
        if shape.arrow_start:
            attrib["marker-start"] = "url(#arrow-tail)"
        if shape.arrow_end:
            attrib["marker-end"] = "url(#arrow-tip)"
        ET.SubElement(parent, f"{{{SVG_NS}}}line", attrib=attrib)
        return
    attrib = {
        "points": point_string(shape.points),
        "fill": "none",
        "stroke": to_svg_color(shape.stroke_color),
        "stroke-width": format_number(max(1.0, shape.stroke_width)),
        "stroke-linecap": "round",
        "stroke-linejoin": "round",
    }
    if shape.arrow_start:
        attrib["marker-start"] = "url(#arrow-tail)"
    if shape.arrow_end:
        attrib["marker-end"] = "url(#arrow-tip)"
    ET.SubElement(parent, f"{{{SVG_NS}}}polyline", attrib=attrib)


def point_string(points: tuple[Point, ...]) -> str:
    return " ".join(f"{format_number(point.x)},{format_number(point.y)}" for point in points)


def bbox_tuple(bbox: BBox) -> tuple[int, int, int, int]:
    return (
        int(round(bbox.x0)),
        int(round(bbox.y0)),
        int(round(bbox.x1)),
        int(round(bbox.y1)),
    )


def flatten_points(points: tuple[Point, ...]) -> tuple[int, ...]:
    flattened: list[int] = []
    for point in points:
        flattened.extend((int(round(point.x)), int(round(point.y))))
    return tuple(flattened)


def draw_arrowhead(
    draw: ImageDraw.ImageDraw,
    start: Point,
    end: Point,
    color: tuple[int, int, int],
    stroke_width: float,
) -> None:
    dx = end.x - start.x
    dy = end.y - start.y
    length = math.hypot(dx, dy)
    if length <= 1.0:
        return
    ux = dx / length
    uy = dy / length
    nx = -uy
    ny = ux
    size = max(7.0, stroke_width * 3.6)
    base = Point(end.x - ux * size, end.y - uy * size)
    left = (int(round(base.x + nx * size * 0.45)), int(round(base.y + ny * size * 0.45)))
    right = (int(round(base.x - nx * size * 0.45)), int(round(base.y - ny * size * 0.45)))
    tip = (int(round(end.x)), int(round(end.y)))
    draw.polygon((tip, left, right), fill=color)
