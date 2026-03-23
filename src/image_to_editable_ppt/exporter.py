from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

from .config import PipelineConfig
from .ir import BoxGeometry, Element, Point, PolylineGeometry


def export_to_pptx(
    elements: list[Element],
    image_size: tuple[int, int],
    output_path: str | Path,
    config: PipelineConfig,
) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    padding = int(Pt(config.slide_padding_pt))
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    image_width, image_height = image_size
    available_width = slide_width - padding * 2
    available_height = slide_height - padding * 2
    scale = min(available_width / image_width, available_height / image_height)
    offset_x = int((slide_width - image_width * scale) / 2.0)
    offset_y = int((slide_height - image_height * scale) / 2.0)

    ordered = sorted(elements, key=z_order_rank)
    for element in ordered:
        add_element(slide, element, scale=scale, offset_x=offset_x, offset_y=offset_y)
    presentation.save(str(output_path))


def z_order_rank(element: Element) -> tuple[int, float]:
    if element.kind in {"rect", "rounded_rect"}:
        return (0, element.bbox.area)
    if element.kind == "text":
        return (2, element.bbox.area)
    return (1, element.bbox.area)


def add_element(slide, element: Element, *, scale: float, offset_x: int, offset_y: int) -> None:
    if element.kind in {"rect", "rounded_rect"}:
        add_box(slide, element, scale=scale, offset_x=offset_x, offset_y=offset_y)
        return
    if element.kind == "line":
        add_line(slide, element, scale=scale, offset_x=offset_x, offset_y=offset_y)
        return
    if element.kind == "orthogonal_connector":
        add_open_polyline(slide, element, scale=scale, offset_x=offset_x, offset_y=offset_y)
        return
    if element.kind == "arrow":
        add_arrow(slide, element, scale=scale, offset_x=offset_x, offset_y=offset_y)
        return
    if element.kind == "text":
        add_textbox(slide, element, scale=scale, offset_x=offset_x, offset_y=offset_y)


def add_box(slide, element: Element, *, scale: float, offset_x: int, offset_y: int) -> None:
    geometry = element.geometry
    if not isinstance(geometry, BoxGeometry):
        raise TypeError("box element requires BoxGeometry")
    bbox = geometry.bbox
    left = to_emu(bbox.x0, scale, offset_x)
    top = to_emu(bbox.y0, scale, offset_y)
    width = max(Emu(1), Emu(int(bbox.width * scale)))
    height = max(Emu(1), Emu(int(bbox.height * scale)))
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if element.kind == "rounded_rect"
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    apply_line_style(shape, element, scale=scale)
    apply_fill_style(shape, element)


def add_line(slide, element: Element, *, scale: float, offset_x: int, offset_y: int) -> None:
    geometry = element.geometry
    if not isinstance(geometry, PolylineGeometry) or len(geometry.points) != 2:
        raise TypeError("line element requires 2-point PolylineGeometry")
    start, end = geometry.points
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        to_emu(start.x, scale, offset_x),
        to_emu(start.y, scale, offset_y),
        to_emu(end.x, scale, offset_x),
        to_emu(end.y, scale, offset_y),
    )
    apply_line_style(shape, element, scale=scale)


def add_open_polyline(slide, element: Element, *, scale: float, offset_x: int, offset_y: int) -> None:
    geometry = element.geometry
    if not isinstance(geometry, PolylineGeometry) or len(geometry.points) < 2:
        raise TypeError("polyline element requires PolylineGeometry")
    start = geometry.points[0]
    builder = slide.shapes.build_freeform(
        start_x=to_emu(start.x, scale, offset_x),
        start_y=to_emu(start.y, scale, offset_y),
    )
    vertices = [
        (to_emu(point.x, scale, offset_x), to_emu(point.y, scale, offset_y))
        for point in geometry.points[1:]
    ]
    builder.add_line_segments(vertices, close=False)
    shape = builder.convert_to_shape()
    apply_line_style(shape, element, scale=scale)
    shape.fill.background()


def add_arrow(slide, element: Element, *, scale: float, offset_x: int, offset_y: int) -> None:
    geometry = element.geometry
    if not isinstance(geometry, PolylineGeometry) or len(geometry.points) != 2:
        raise TypeError("arrow element requires 2-point PolylineGeometry")
    if try_add_semantic_arrow(slide, element, scale=scale, offset_x=offset_x, offset_y=offset_y):
        return
    # python-pptx does not expose arrowheads directly, so freeform is the conservative fallback.
    polygon = arrow_polygon(geometry.points[0], geometry.points[1], max(4.0, element.stroke.width * 1.4))
    builder = slide.shapes.build_freeform(
        start_x=to_emu(polygon[0].x, scale, offset_x),
        start_y=to_emu(polygon[0].y, scale, offset_y),
    )
    vertices = [
        (to_emu(point.x, scale, offset_x), to_emu(point.y, scale, offset_y))
        for point in polygon[1:]
    ]
    builder.add_line_segments(vertices, close=True)
    shape = builder.convert_to_shape()
    rgb = RGBColor(*element.stroke.color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.color.rgb = rgb
    shape.line.width = Emu(max(1, int(element.stroke.width * scale * 0.6)))


def try_add_semantic_arrow(slide, element: Element, *, scale: float, offset_x: int, offset_y: int) -> bool:
    geometry = element.geometry
    if not isinstance(geometry, PolylineGeometry) or len(geometry.points) != 2:
        return False
    start, end = geometry.points
    try:
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            to_emu(start.x, scale, offset_x),
            to_emu(start.y, scale, offset_y),
            to_emu(end.x, scale, offset_x),
            to_emu(end.y, scale, offset_y),
        )
        apply_line_style(shape, element, scale=scale)
        # DrawingML names the connector start as "head" and the connector end as "tail".
        # Arrow geometry is normalized so geometry.points[1] is the arrow tip, which maps to
        # the connector end and therefore needs a:tailEnd markup.
        add_connector_arrowhead(shape, end_at="tail")
        return True
    except Exception:
        return False


def add_textbox(slide, element: Element, *, scale: float, offset_x: int, offset_y: int) -> None:
    geometry = element.geometry
    if not isinstance(geometry, BoxGeometry):
        raise TypeError("text element requires BoxGeometry")
    text = "" if element.text is None else element.text.content
    bbox = geometry.bbox
    shape = slide.shapes.add_textbox(
        to_emu(bbox.x0, scale, offset_x),
        to_emu(bbox.y0, scale, offset_y),
        Emu(max(1, int(bbox.width * scale))),
        Emu(max(1, int(bbox.height * scale))),
    )
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = alignment_for_element(element)
    shape.fill.background()
    shape.line.fill.background()


def alignment_for_element(element: Element):
    alignment = "center" if element.text is None else element.text.alignment
    if alignment == "left":
        return PP_ALIGN.LEFT
    if alignment == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.CENTER


def apply_line_style(shape, element: Element, *, scale: float) -> None:
    rgb = RGBColor(*element.stroke.color)
    shape.line.color.rgb = rgb
    shape.line.width = Emu(max(1, int(max(1.0, element.stroke.width) * scale)))
    if element.stroke.dash_style != "solid":
        try:
            shape.line.dash_style = element.stroke.dash_style
        except Exception:
            pass


def apply_fill_style(shape, element: Element) -> None:
    if element.fill.enabled and element.fill.color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*element.fill.color)
    else:
        shape.fill.background()


def arrow_polygon(start: Point, end: Point, width_pixels: float) -> tuple[Point, ...]:
    dx = end.x - start.x
    dy = end.y - start.y
    length = math.hypot(dx, dy)
    if length <= 1.0:
        return (start, end, end)
    ux = dx / length
    uy = dy / length
    nx = -uy
    ny = ux
    shaft_half = max(2.5, width_pixels * 0.9)
    head_half = shaft_half * 2.0
    head_length = min(length * 0.45, max(8.0, shaft_half * 4.0))
    body_length = max(length - head_length, shaft_half * 2.0)
    tail_left = Point(start.x + nx * shaft_half, start.y + ny * shaft_half)
    tail_right = Point(start.x - nx * shaft_half, start.y - ny * shaft_half)
    neck_left = Point(start.x + ux * body_length + nx * shaft_half, start.y + uy * body_length + ny * shaft_half)
    neck_right = Point(start.x + ux * body_length - nx * shaft_half, start.y + uy * body_length - ny * shaft_half)
    head_left = Point(start.x + ux * body_length + nx * head_half, start.y + uy * body_length + ny * head_half)
    head_right = Point(start.x + ux * body_length - nx * head_half, start.y + uy * body_length - ny * head_half)
    tip = end
    return (tail_left, tail_right, neck_right, head_right, tip, head_left, neck_left)


def to_emu(value: float, scale: float, offset: int) -> Emu:
    return Emu(int(offset + value * scale))


def add_connector_arrowhead(shape, *, end_at: str) -> None:
    ln = shape._element.spPr.get_or_add_ln()
    tag = qn("a:tailEnd" if end_at == "tail" else "a:headEnd")
    for child in list(ln):
        if child.tag == tag:
            ln.remove(child)
    arrow = OxmlElement("a:tailEnd" if end_at == "tail" else "a:headEnd")
    arrow.set("type", "triangle")
    arrow.set("w", "med")
    arrow.set("len", "med")
    ln.append(arrow)
