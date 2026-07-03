"""Minimal editable-pptx writer for an EmitScene (Phase 10).

Writes native PowerPoint primitives — autoshapes, text boxes, connectors — so
the user can immediately move / edit / delete everything. Coordinate contract
matches the synthetic generator: 1 px = 9525 EMU at 96 dpi, slide size derived
from the image size, so the pptx mirrors the (possibly resized) input image.

Conservative rules (principle: 확실한 것만 완벽하게):
- connectors without path geometry are not drawn;
- standalone text regions without recognized text are not drawn (a blank is
  honest, a hallucinated label is not);
- text owned by a node goes into the node's text frame so it moves with the
  shape; everything else becomes a borderless text box.
"""
from __future__ import annotations

from pathlib import Path
from typing import Mapping

from image_to_editable_ppt.v3.core.enums import NodeKind, PortOwnerKind
from image_to_editable_ppt.v3.emit.models import (
    EmitConnectorPrimitive,
    EmitScene,
    EmitShapePrimitive,
    EmitTextPrimitive,
)
from image_to_editable_ppt.v3.emit.style import ShapeVisualStyle

EMU_PER_PIXEL = 9525  # 1 px at 96 dpi; same contract as ml/synthesize.py
_TEXT_COLOR = (15, 23, 42)
_CONNECTOR_COLOR = (31, 41, 55)
_DEFAULT_STYLE = ShapeVisualStyle(stroke=(71, 85, 105))


def write_pptx(
    scene: EmitScene,
    path: str | Path,
    *,
    styles: Mapping[str, ShapeVisualStyle] | None = None,
) -> None:
    from pptx import Presentation
    from pptx.util import Emu

    output_path = Path(path)
    presentation = Presentation()
    presentation.slide_width = Emu(scene.image_size.width * EMU_PER_PIXEL)
    presentation.slide_height = Emu(scene.image_size.height * EMU_PER_PIXEL)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    owned_text_ids: set[str] = set()
    node_shape_ids = {shape.id for shape in scene.shapes if shape.owner_kind is PortOwnerKind.NODE}
    texts_by_owner: dict[str, list[EmitTextPrimitive]] = {}
    for text in scene.texts:
        for owner_id in text.owner_ids:
            if owner_id in node_shape_ids and text.text:
                texts_by_owner.setdefault(owner_id, []).append(text)
                owned_text_ids.add(text.id)
                break

    # z-order: containers behind nodes, larger shapes behind smaller ones.
    containers = sorted(
        (shape for shape in scene.shapes if shape.owner_kind is PortOwnerKind.CONTAINER),
        key=lambda shape: -shape.bbox.area,
    )
    nodes = sorted(
        (shape for shape in scene.shapes if shape.owner_kind is PortOwnerKind.NODE),
        key=lambda shape: -shape.bbox.area,
    )
    for shape in (*containers, *nodes):
        style = (styles or {}).get(shape.id, _DEFAULT_STYLE)
        _add_shape(slide, shape, style=style, owned_texts=texts_by_owner.get(shape.id, ()))

    for connector in scene.connectors:
        _add_connector(slide, connector)

    for text in scene.texts:
        if text.id in owned_text_ids or not text.text:
            continue
        _add_textbox(slide, text)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def _emu_box(bbox) -> tuple:
    from pptx.util import Emu

    return (
        Emu(int(bbox.x0 * EMU_PER_PIXEL)),
        Emu(int(bbox.y0 * EMU_PER_PIXEL)),
        Emu(int(max(1.0, bbox.width) * EMU_PER_PIXEL)),
        Emu(int(max(1.0, bbox.height) * EMU_PER_PIXEL)),
    )


def _add_shape(
    slide,
    shape: EmitShapePrimitive,
    *,
    style: ShapeVisualStyle,
    owned_texts: tuple[EmitTextPrimitive, ...] | list[EmitTextPrimitive],
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    label_text, text_extent = _merge_owned_texts(owned_texts, fallback_label=shape.label)

    if shape.shape_kind is NodeKind.LABEL_ANCHOR:
        # Text-only node (e.g. tree label): a borderless text box at the node extent.
        box = slide.shapes.add_textbox(*_emu_box(shape.bbox))
        box.name = shape.id
        box.fill.background()
        box.line.fill.background()
        if label_text:
            _fill_text_frame(box.text_frame, label_text, extent=text_extent or (shape.bbox.width, shape.bbox.height), wrap=False)
        return

    mso_shape = MSO_SHAPE.ROUNDED_RECTANGLE if shape.shape_kind is NodeKind.ROUNDED_BOX else MSO_SHAPE.RECTANGLE
    autoshape = slide.shapes.add_shape(mso_shape, *_emu_box(shape.bbox))
    autoshape.name = shape.id
    autoshape.shadow.inherit = False
    if style.fill is not None:
        autoshape.fill.solid()
        autoshape.fill.fore_color.rgb = RGBColor(*style.fill)
    else:
        autoshape.fill.background()
    if style.stroke is not None:
        from pptx.util import Pt

        autoshape.line.color.rgb = RGBColor(*style.stroke)
        autoshape.line.width = Pt(max(0.75, style.stroke_width * 0.75))
    else:
        autoshape.line.fill.background()
    if label_text:
        _fill_text_frame(
            autoshape.text_frame,
            label_text,
            extent=text_extent or (shape.bbox.width * 0.8, shape.bbox.height * 0.4),
            wrap=True,
        )


def _merge_owned_texts(
    owned_texts: tuple[EmitTextPrimitive, ...] | list[EmitTextPrimitive],
    *,
    fallback_label: str | None,
) -> tuple[str | None, tuple[float, float] | None]:
    """Join owned text regions top-to-bottom; also return their union extent (px)."""
    if not owned_texts:
        return fallback_label, None
    ordered = sorted(owned_texts, key=lambda text: (text.bbox.y0, text.bbox.x0))
    merged = "\n".join(text.text for text in ordered if text.text)
    width = max(text.bbox.x1 for text in ordered) - min(text.bbox.x0 for text in ordered)
    height = max(text.bbox.y1 for text in ordered) - min(text.bbox.y0 for text in ordered)
    return (merged or fallback_label), (width, height)


def _font_size_pt(text: str, *, extent: tuple[float, float]) -> int:
    """Font size that keeps the text inside its observed pixel extent.

    Height bound: line box ≈ 1.3 em at 96 dpi (0.55 pt/px after the region's
    padding). Width bound: average glyph advance ≈ 0.5 em, so the longest line
    must fit `1.4 * width / chars` pt. Without the width bound, long single-line
    regions overflow their box and overlap neighbours.
    """
    lines = [line for line in text.split("\n") if line] or [text]
    longest = max(len(line) for line in lines)
    width_px, height_px = extent
    height_bound = 0.55 * (height_px / len(lines))
    width_bound = 1.4 * (width_px / max(1, longest))
    return max(6, min(24, round(min(height_bound, width_bound))))


def _fill_text_frame(text_frame, text: str, *, extent: tuple[float, float], wrap: bool) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Emu, Pt

    text_frame.word_wrap = wrap
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = Emu(0)
    text_frame.margin_right = Emu(0)
    text_frame.margin_top = Emu(0)
    text_frame.margin_bottom = Emu(0)
    font_size = Pt(_font_size_pt(text, extent=extent))
    for index, line in enumerate(text.split("\n")):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = font_size
            run.font.color.rgb = RGBColor(*_TEXT_COLOR)


def _add_connector(slide, connector: EmitConnectorPrimitive) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Pt

    points = connector.path_points
    if len(points) < 2:
        return  # no geometry, nothing honest to draw
    for start, end in zip(points, points[1:]):
        segment = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(int(start.x * EMU_PER_PIXEL)),
            Emu(int(start.y * EMU_PER_PIXEL)),
            Emu(int(end.x * EMU_PER_PIXEL)),
            Emu(int(end.y * EMU_PER_PIXEL)),
        )
        segment.name = connector.id
        segment.line.color.rgb = RGBColor(*_CONNECTOR_COLOR)
        segment.line.width = Pt(1.25)
        line_element = segment.line._get_or_add_ln()
        if connector.arrowhead_end and end is points[-1]:
            line_element.append(line_element.makeelement(qn("a:tailEnd"), {"type": "triangle"}))
        if connector.arrowhead_start and start is points[0]:
            line_element.append(line_element.makeelement(qn("a:headEnd"), {"type": "triangle"}))


def _add_textbox(slide, text: EmitTextPrimitive) -> None:
    box = slide.shapes.add_textbox(*_emu_box(text.bbox))
    box.name = text.id
    box.fill.background()
    box.line.fill.background()
    _fill_text_frame(box.text_frame, text.text or "", extent=(text.bbox.width, text.bbox.height), wrap=False)
